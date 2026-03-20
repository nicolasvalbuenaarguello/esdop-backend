from flask import make_response
from datetime import date, time, datetime
from tipo_docker.b_f_resultados_lineas_estrategicas_paya_power.models.estadistica.estadistica_boletin_coe import *
from __init__ import *
 
from flask import make_response
from datetime import date, time, datetime



from tipo_docker.z_z_configuarcion.caligrafia import *
from tipo_docker.z_z_configuarcion.fechas import *
from tipo_docker.z_z_configuarcion.header import *
from tipo_docker.z_z_configuarcion.logo import *
from tipo_docker.z_z_configuarcion.titulos import *  
from tipo_docker.z_z_configuarcion.tipo_ayudas import *

from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.enum.chart import XL_TICK_MARK
from pptx.util import Pt

from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.chart import XL_TICK_LABEL_POSITION

from pptx.enum.shapes import MSO_SHAPE

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_VERTICAL_ANCHOR
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.dml import MSO_COLOR_TYPE
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_FILL
from pptx.oxml.xmlchemy import OxmlElement
import os
import shutil

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL



from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL


from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

import matplotlib.pyplot as plt
import numpy as np
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.dml.color import RGBColor



def grafico_reloj(ruta_salida, velocidad=0, color_reloj=(0,176,80), nombre='dirop'):
        """
        Genera un velocímetro tipo reloj y guarda la imagen PNG en la ruta indicada.
        Parámetros:
        - ruta_salida: ruta completa del archivo a guardar (sin extensión o con .png)
        - velocidad: valor entre 0 y 100
        """

        velocidad_max = 100
        velocidad = max(0, min(velocidad, velocidad_max))  # límite 0–100

        # --- Cálculo del ángulo ---
        angulo = np.deg2rad(-145 + (velocidad / velocidad_max) * 290)

        # --- Crear figura polar ---
        fig, ax = plt.subplots(figsize=(6, 6), subplot_kw={'projection': 'polar'})
        ax.set_theta_zero_location('N')
        ax.set_theta_direction(-1)
        ax.set_thetamin(-145)
        ax.set_thetamax(145)
        ax.set_ylim(0, 1)
        ax.set_yticklabels([])
        ax.set_xticklabels([])
        ax.grid(False)
        ax.spines['polar'].set_visible(False)

        # --- Fondo gris del arco ---
        theta_full = np.linspace(np.deg2rad(-145), np.deg2rad(145), 500)
        r_inner, r_outer = 0.1, 0.9
        for i in range(len(theta_full) - 1):
            ax.fill_between(
                [theta_full[i], theta_full[i + 1]],
                r_inner, r_outer,
                color="#D0D0D0",
                alpha=0.6
            )

        # --- Colorear según la velocidad ---
        def color_por_velocidad(v):
            if v < 35:
                return "#FF0000"     # rojo
            elif v < 50:
                return "#FFA600"     # naranja
            elif v < 80:
                return "#FFFF00"     # amarillo
            else:
                return "#00FF00"     # verde

        color = color_reloj

        theta_color = np.linspace(
            np.deg2rad(-145),
            np.deg2rad(-145 + (velocidad / velocidad_max) * 290),
            300
        )

        for i in range(len(theta_color) - 1):
            ax.fill_between(
                [theta_color[i], theta_color[i + 1]],
                r_inner, r_outer,
                color=color,
                alpha=0.9
            )

        # --- Líneas divisorias y numeración ---
        for v in range(0, velocidad_max + 1, 10):
            ang = np.deg2rad(-145 + (v / velocidad_max) * 290)
            ax.plot([ang, ang], [r_inner, r_outer], color='white', linewidth=2, zorder=6)

            # Texto de numeración
            r_text = r_outer + 0.08  # posición radial del texto
            ax.text(ang, r_text, str(v), ha='center', va='center',
                    color='gray', fontsize=25, fontweight='bold')

        # --- Aguja ---
        ax.plot([angulo, angulo], [0, 0.85], color='black', linewidth=3, zorder=10)
        ax.scatter(0, 0, color='black', s=120, zorder=11)


        # --- Fondo y guardado ---
        fig.patch.set_facecolor("#00000000")
        ax.set_facecolor("#404040")

        # Asegurar extensión PNG con nombre fijo
        nombre = nombre + '.png'
        if not ruta_salida.lower().endswith(nombre):
            if not ruta_salida.endswith(("/", "\\")):
                ruta_salida += "_"
            ruta_salida += nombre

        plt.savefig(ruta_salida, dpi=100, bbox_inches='tight', transparent=True)
        plt.close(fig)

        return ruta_salida

def set_cell_border(cell):
    """Aplica bordes negros a una celda de tabla en pptx."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for side in ['L', 'R', 'T', 'B']:
        ln = OxmlElement(f'a:ln{side}')
        ln.set('w', '12700')  # 1pt en EMUs

        solidFill = OxmlElement('a:solidFill')
        srgbClr = OxmlElement('a:srgbClr')
        srgbClr.set('val', '000000')
        solidFill.append(srgbClr)
        ln.append(solidFill)

        prstDash = OxmlElement('a:prstDash')
        prstDash.set('val', 'solid')
        ln.append(prstDash)

        tcPr.append(ln)

def apply_border(cell, color="000000", width=1):

    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    for line in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:

        ln = OxmlElement(line)
        ln.set('w', str(width * 12700))

        solidFill = OxmlElement('a:solidFill')
        srgbClr = OxmlElement('a:srgbClr')
        srgbClr.set('val', color)

        solidFill.append(srgbClr)
        ln.append(solidFill)

        tcPr.append(ln)


def compactar_texto_celda(cell, texto, font_name="Calibri", size_pt=9, bold=True, align=PP_ALIGN.CENTER):

    tf = cell.text_frame
    tf.clear()

    # quitar padding
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.margin_left = 0
    tf.margin_right = 0

    tf.word_wrap = False
    tf.auto_size = None
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.text = texto

    # 🔹 eliminar espacios extra
    p.space_before = Pt(0)
    p.space_after = Pt(0)

    # 🔹 interlineado simple
    p.line_spacing = 0.7

    # alineación horizontal
    p.alignment = align

    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold

def crear_tabla_afectacion(
        prs,
        slide,
        data,
        size_pt=9,
        left=0.5,
        top=1.2,
        width=9,
        font_name="Calibri",
        row_height_cm=0.5,
        header_height_cm=None,
        col_widths = [
        Inches(0.45),
        Inches(1.2),
        Inches(3),
        Inches(1),
        Inches(1),
        Inches(1.5),
        Inches(1.5),
        Inches(1.5)
        ],
        header_color = RGBColor(127,127,127),
        indicadore = 0
):

    rows = len(data)
    cols = len(data[0])

    # ------------------------------------------------
    # ALTURA FILAS
    # ------------------------------------------------
    row_height = Inches(row_height_cm / 2.54)
    header_height = Inches(header_height_cm / 2.54) if header_height_cm else row_height

    total_height = row_height * rows

    # ------------------------------------------------
    # CREAR TABLA
    # ------------------------------------------------
    table_shape = slide.shapes.add_table(
        rows,
        cols,
        Inches(left),
        Inches(top),
        Inches(width),
        total_height
    )

    table = table_shape.table

    # ------------------------------------------------
    # ALTURA FILAS
    # ------------------------------------------------
    for i in range(rows):
        table.rows[i].height = header_height if i == 0 else row_height

    # ------------------------------------------------
    # ANCHO COLUMNAS
    # ------------------------------------------------

    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # ------------------------------------------------
    # COLORES
    # ------------------------------------------------
    header_color = header_color
    white = RGBColor(255,255,255)
    black = RGBColor(0,0,0)

    fila_clara = RGBColor(217,217,217)
    fila_blanca = RGBColor(242,242,242)

    rojo = RGBColor(192,0,0)
    verde = RGBColor(0,176,80)
    amarillo = RGBColor(255,192,0)

    # ------------------------------------------------
    # RECORRER DATOS
    # ------------------------------------------------
    for i, row in enumerate(data):

        for j, val in enumerate(row):

            cell = table.cell(i, j)

            # ----------------------------------------
            # ALINEACIÓN
            # ----------------------------------------
            if indicadore == 0:
                if i == 0:
                    alineacion = PP_ALIGN.CENTER
                elif j == 2:
                    alineacion = PP_ALIGN.LEFT
                else:
                    alineacion = PP_ALIGN.CENTER
            else:
                if i == 0:
                    alineacion = PP_ALIGN.CENTER
                elif j == 0:
                    alineacion = PP_ALIGN.LEFT
                else:
                    alineacion = PP_ALIGN.CENTER


            # ----------------------------------------
            # TEXTO COMPACTO
            # ----------------------------------------
            compactar_texto_celda(
                cell,
                str(val),
                font_name=font_name,
                size_pt=size_pt,
                bold=True,
                align=alineacion
            )

            run = cell.text_frame.paragraphs[0].runs[0]

            # ----------------------------------------
            # ENCABEZADO
            # ----------------------------------------
            if i == 0:

                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
                run.font.color.rgb = white

            else:

                # ----------------------------------------
                # FILAS ALTERNAS
                # ----------------------------------------
                cell.fill.solid()

                if i % 2 == 0:
                    cell.fill.fore_color.rgb = fila_clara
                else:
                    cell.fill.fore_color.rgb = fila_blanca

                run.font.color.rgb = black

                # ----------------------------------------
                # SEMÁFORO
                # ----------------------------------------
                if j == 7:

                    valor = str(val).strip()

                    if valor == "+":
                        cell.fill.fore_color.rgb = verde
                        run.font.color.rgb = white

                    elif valor == "-":
                        cell.fill.fore_color.rgb = rojo
                        run.font.color.rgb = white

                    elif valor == "=":
                        cell.fill.fore_color.rgb = amarillo
                        run.font.color.rgb = black

                if j == 6:

                    valor = str(val).strip()

                    if "-" == valor:
                        run.font.color.rgb = black
                    elif "-" in valor:
                        run.font.color.rgb = rojo
                        
                if j == 5:

                    valor = str(val).strip()

                    if "-" == valor:
                        run.font.color.rgb = black 
                    elif "-" in valor:
                        run.font.color.rgb = rojo


            # ----------------------------------------
            # BORDES
            # ----------------------------------------
            apply_border(cell)

    return table

def cuadros_divisio(slide, valor_division, x, y):
  # ... dentro de tu función o flujo
    shapes = slide.shapes
    left = Inches(x)
    top = Inches(y)

    width = Inches(1.1)
    height = Inches(0.3)

    shape = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )

    # Eliminar fondo del rectángulo
    shape.fill.background()

    # Eliminar bordes
    shape.line.fill.background()

    # Agregar texto
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = str(valor_division)

    # Estilo de texto
    font = run.font
    font.name = 'Arial'
    font.size = Pt(14)
    font.bold = True
    font.color.rgb = RGBColor(0, 0, 0)

    # Alineación opcional
    p.alignment = PP_ALIGN.LEFT

def cuadros_divisio_r(slide, valor_division, x, y, rotar, color, ancho):
    shapes = slide.shapes
    left = Inches(x)
    top = Inches(y)
    width = Inches(ancho)
    height = Inches(0.3)

    shape = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )

    # Eliminar fondo y bordes
    shape.fill.background()
    shape.line.fill.background()

    # Agregar texto
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = str(valor_division)

    # Estilo de texto
    font = run.font
    font.name = color[3]
    font.size = Pt(color[4])
    font.bold = color[5]
    font.italic = color[6]
    font.color.rgb = RGBColor(color[0],color[1],color[2])

    # Alineación
    p.alignment = PP_ALIGN.CENTER

    # 🔄 Rotar el texto o el shape (90 grados)
    shape.rotation = rotar  # rota todo el cuadro de texto

    # Si quieres solo el texto en vertical sin girar la forma:
    # text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    # text_frame.text_anchor = MSO_VERTICAL_ANCHOR.VERTICAL

# Función para agregar línea personalizada
def agregar_linea(slide, left_in, top_in, largo_in, orientacion="horizontal", color_rgb=(0, 0, 0), grosor_pt=2):
    left = Inches(left_in)
    top = Inches(top_in)

    if orientacion == "horizontal":
        width = Inches(largo_in)
        height = Inches(0)  # Línea horizontal
    elif orientacion == "vertical":
        width = Inches(0)
        height = Inches(largo_in)  # Línea vertical
    elif orientacion == "diagonal":
        width = Inches(largo_in)
        height = Inches(largo_in)  # Diagonal con 45°
    else:
        raise ValueError("Orientación no válida. Usa: horizontal, vertical o diagonal.")

    line = slide.shapes.add_shape(
        autoshape_type_id=1,  # Línea
        left=left,
        top=top,
        width=width,
        height=height
    )

    # Estilizar línea
    line.line.color.rgb = RGBColor(*color_rgb)
    line.line.width = Pt(grosor_pt)

def merge_repetidos_columna(tabla, data, col_index=1,
                            font_name="Calibri",
                            font_size=11,
                            color=(0,0,0),
                            bold=True):

    # -----------------------------
    # VALIDAR DATA
    # -----------------------------
    if not data:
        return

    # eliminar filas vacías
    data = [row for row in data if row]

    # si después de limpiar queda menos de 2 filas
    if len(data) < 2:
        return

    # verificar columnas
    for row in data:
        if len(row) <= col_index:
            return

    start = 1
    actual = data[start][col_index]

    for i in range(start + 1, len(data)):

        if data[i][col_index] != actual:

            end = i - 1

            if end > start:
                merge_clean(
                    tabla,
                    start, col_index,
                    end, col_index,
                    text=actual,
                    font_name=font_name,
                    font_size=font_size,
                    color=color,
                    bold=bold
                )

            start = i
            actual = data[i][col_index]

    # -----------------------------
    # último grupo
    # -----------------------------
    end = len(data) - 1
    if end > start:
        merge_clean(
            tabla,
            start, col_index,
            end, col_index,
            text=actual,
            font_name=font_name,
            font_size=font_size,
            color=color,
            bold=bold
        )

def merge_clean(
        table,
        r1, c1,
        r2, c2,
        text=None,
        font_name="Calibri",
        font_size=9,
        color=(0,0,0),
        bold=True,
        align="center"
    ):

    # ----------------------------------------
    # limpiar celdas
    # ----------------------------------------
    for r in range(r1, r2+1):
        for c in range(c1, c2+1):
            if not (r == r1 and c == c1):
                table.cell(r,c).text = ""

    # ----------------------------------------
    # fusionar
    # ----------------------------------------
    try:
        table.cell(r1,c1).merge(table.cell(r2,c2))
    except:
        pass

    cell = table.cell(r1,c1)

    # ----------------------------------------
    # centrado vertical
    # ----------------------------------------
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcPr.set("anchor", "ctr")

    # ----------------------------------------
    # limpiar text frame
    # ----------------------------------------
    tf = cell.text_frame
    tf.clear()

    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.margin_left = 0
    tf.margin_right = 0

    # ----------------------------------------
    # texto
    # ----------------------------------------
    p = tf.paragraphs[0]
    p.text = text if text else ""

    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1

    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT

    # ----------------------------------------
    # fuente
    # ----------------------------------------
    if not p.runs:
        run = p.add_run()
    else:
        run = p.runs[0]

    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)

    return cell

def agregar_texto(
    slide,
    texto,
    left_in,
    top_in,
    width_in=0.5,
    height_in=0.3,
    font_size=12,
    font_name="Calibri",
    color=(127,127,127),
    bold=False,
    italic=False,
    align="left",
    word_wrap=False,

    # 🔹 NUEVOS PARÁMETROS
    fill=False,
    fill_color=(255,255,255),

    margin_left=0,
    margin_right=0,
    margin_top=0,
    margin_bottom=0,
    underline=False, # 👈 subrayado
    strike=False, # 👈 tachado
    shadow=False
):
    """
    Inserta un textbox en PowerPoint.

    Por defecto:
    - Sin relleno
    - Sin márgenes
    - Texto alineado según parámetro

    Parámetros nuevos:
    fill -> activar relleno
    fill_color -> color del relleno
    margin_* -> márgenes internos
    """

    # -------------------------------------------------
    # POSICIÓN Y TAMAÑO
    # -------------------------------------------------
    left = Inches(left_in)
    top = Inches(top_in)
    width = Inches(width_in)
    height = Inches(height_in)

    textbox = slide.shapes.add_textbox(left, top, width, height)

    # -------------------------------------------------
    # CONFIGURAR RELLENO
    # -------------------------------------------------
    if fill:
        textbox.fill.solid()
        textbox.fill.fore_color.rgb = RGBColor(*fill_color)
    else:
        textbox.fill.background()   # sin relleno

    # -------------------------------------------------
    # TEXTO
    # -------------------------------------------------
    tf = textbox.text_frame

    tf.word_wrap = word_wrap
    tf.auto_size = None
    tf.clear()

    # -------------------------------------------------
    # MÁRGENES
    # -------------------------------------------------
    tf.margin_left = Pt(margin_left)
    tf.margin_right = Pt(margin_right)
    tf.margin_top = Pt(margin_top)
    tf.margin_bottom = Pt(margin_bottom)

    # -------------------------------------------------
    # PÁRRAFO
    # -------------------------------------------------
    p = tf.paragraphs[0]
    p.text = texto

    # alineación
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    elif align == "justify":
        p.alignment = PP_ALIGN.JUSTIFY
    else:
        p.alignment = PP_ALIGN.LEFT

    # -------------------------------------------------
    # FUENTE
    # -------------------------------------------------
    font = p.font
    font.size = Pt(font_size)
    font.name = font_name
    font.bold = bold
    font.italic = italic
    font.color.rgb = RGBColor(*color)
    font.underline = underline   # 👈 subrayado
    font.strike = strike         # 👈 tachado
    font.shadow = shadow

    return textbox



def agregar_cuadro_redondo(
    slide,
    texto,
    left_in,
    top_in,
    width_in=0.5,
    height_in=0.3,
    font_size=12,
    font_name="Calibri",
    color=(0,0,0),

    bold=False,
    italic=False,
    align="left",

    fill=True,
    fill_color=(255,255,255),

    border=True,
    border_color=(0,0,0),
    border_width=1,

    margin_left=5,
    margin_right=5,
    margin_top=3,
    margin_bottom=3
):
    
    # -------------------------------------------------
    # POSICIÓN
    # -------------------------------------------------
    left = Inches(left_in)
    top = Inches(top_in)
    width = Inches(width_in)
    height = Inches(height_in)

    # -------------------------------------------------
    # CREAR RECTÁNGULO REDONDO
    # -------------------------------------------------
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height
    )

    # -------------------------------------------------
    # RELLENO
    # -------------------------------------------------
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill_color)
    else:
        shape.fill.background()

    # -------------------------------------------------
    # BORDE
    # -------------------------------------------------
    if border:
        shape.line.color.rgb = RGBColor(*border_color)
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()

    # -------------------------------------------------
    # TEXTO
    # -------------------------------------------------
    tf = shape.text_frame
    tf.clear()

    tf.margin_left = Pt(margin_left)
    tf.margin_right = Pt(margin_right)
    tf.margin_top = Pt(margin_top)
    tf.margin_bottom = Pt(margin_bottom)

    p = tf.paragraphs[0]
    p.text = texto

    # alineación
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    elif align == "justify":
        p.alignment = PP_ALIGN.JUSTIFY
    else:
        p.alignment = PP_ALIGN.LEFT

    # -------------------------------------------------
    # FUENTE
    # -------------------------------------------------
    font = p.font
    font.size = Pt(font_size)
    font.name = font_name
    font.bold = bold
    font.italic = italic
    font.color.rgb = RGBColor(*color)

    return shape

def colorear_celdas_unidas(tabla, data, col_index=1):

    # validar que exista contenido
    if not data or len(data) < 2:
        return

    # validar columnas
    for row in data:
        if len(row) <= col_index:
            return

    color1 = RGBColor(242,242,242)
    color2 = RGBColor(217,217,217)

    color_actual = color1
    start = 1
    actual = data[start][col_index]

    for i in range(2, len(data)):

        if data[i][col_index] != actual:

            cell = tabla.cell(start, col_index)
            cell.fill.solid()
            cell.fill.fore_color.rgb = color_actual

            # alternar color
            color_actual = color2 if color_actual == color1 else color1

            start = i
            actual = data[i][col_index]

    # último grupo
    cell = tabla.cell(start, col_index)
    cell.fill.solid()
    cell.fill.fore_color.rgb = color_actual

def reloj_grafico(slide_cuerpo_5, reloj_base, altura, izquierda, img_path, texto, colo_texto, color_relleno):
    slide_cuerpo_5.shapes.add_picture(
    reloj_base,
    Inches(altura),   # izquierda
    Inches(izquierda),     # arriba
    width=Inches(3),
    height=Inches(3)
    )
    altura_1 = altura + 0.6
    izquierda_1 = izquierda + 0.52
    slide_cuerpo_5.shapes.add_picture(
    img_path,
    Inches(altura_1),   # izquierda
    Inches(izquierda_1),     # arriba
    width=Inches(1.8),
    height=Inches(1.8)
    )
    altura = altura + 0.9 
    izquierda = izquierda + 2.6 

    agregar_cuadro_redondo(
        slide_cuerpo_5,
        texto=texto,
        left_in=altura,
        top_in=izquierda,
        width_in=1.3,
        height_in=0.4,
        font_size=14,
        bold=True,
        align="center",
        fill=True,
        color=colo_texto,
        fill_color=color_relleno,
        border_color=color_relleno,
        border_width=2
    )

#comparativo por divisiones
def comparativo_comparativo_mapa(fecha_inicial_u_l, fecha_final_u_l, fecha_inicial_p_l, fecha_final_p_l, filtro, link, puerto):



    
    fechas_inicial = fecha(fecha_inicial_u_l)
    
    fecha_inicial_dia_inicial = fechas_inicial[0]
    fecha_inicial_mes_inicial = fechas_inicial[1]
    fecha_inicial_anio_inicial = fechas_inicial[2]

    fechas = fecha(fecha_final_u_l)

    fecha_inicial_dia = fechas[0]
    fecha_inicial_mes = fechas[1]
    fecha_inicial_anio = fechas[2]
    fecha_inicial_mes_number = fechas[3]

    fechas_inicial = fecha(fecha_inicial_u_l)
    fechas_final = fecha(fecha_final_u_l)

    fechas_inicial_final = fecha(fecha_inicial_p_l)
    fechas_final_final = fecha(fecha_final_p_l)
    
    fecha_titulo = fechas_inicial[0]+" " +  fechas_inicial[1] + " AL " + fechas_final[0] + " "+fechas_final[1] + " "+fechas_final[2]
    fecha_titulo = fecha_titulo.upper()

    fecha_titulo_dos = fechas_inicial_final[0]+" " +  fechas_inicial_final[1] + " AL " + fechas_final_final[0] + " "+fechas_final_final[1] + " "+fechas_final_final[2]
    fecha_titulo_dos = fecha_titulo_dos.upper()



    if fecha_inicial_anio_inicial == fecha_inicial_anio:
        fechas= str("EL ")+str(fecha_inicial_dia_inicial) + " DE " +str(fecha_inicial_mes_inicial) + " hasta " + str(fecha_inicial_dia) + " de "+ str(fecha_inicial_mes) + " del " + str(fechas_final_final[2])+ " - " +str(fecha_inicial_anio)
    else:
        fechas= str("EL ")+str(fecha_inicial_dia_inicial) + " DE " +str(fecha_inicial_mes_inicial) + " DEL " +str(fecha_inicial_anio_inicial) + " hasta " + str(fecha_inicial_dia) + " de "+ str(fecha_inicial_mes) + " del " + str(fecha_inicial_anio)
    fechas = fechas.upper()

    

    unidad = ""
    if filtro[0] != "---" and filtro[0] != "":
        unidad = filtro[0]

    titulo="JEFATURA DE ESTADO MAYOR DE OPERACIONES {}".format(unidad)
    fecha_titulo_sub = 'RESULTADOS INTERDICION COMPARATIVO {}-{}'.format(fechas_final[2],fechas_final_final[2])

    resultados_spoa =  Calculo_Spoa(fecha_inicial_u_l, fecha_final_u_l, fecha_inicial_p_l, fecha_final_p_l, filtro, fechas_final[2], fechas_final_final[2], fecha_titulo, fecha_titulo_dos)
    datos_resultados = resultados_spoa.comparativo_mapa()

    ruta_base = filtro[15]
    ayuda = f"{ruta_base}/static/template/template_docna_paya.pptx"
    prs = Presentation(ayuda)
    #este codigo es para saber que tipo de ayudas tengo
    #tipo_ayudas_template(prs)
   

    slide_portada = prs.slides.add_slide(prs.slide_layouts[1])
    slide_portada = prs.slides.add_slide(prs.slide_layouts[2])
    slide_portada = prs.slides.add_slide(prs.slide_layouts[4])
 

    slide_cuerpo = prs.slides.add_slide(prs.slide_layouts[5])

    titulo = f"""LÍNEA DE OPERACIONES “PAYA” ESTRATEGIA DINAMIZADORA """
    titulo_2 = f"""RESULTADOS DE INTERDICCIÓN DESDE {fechas}  """
    agregar_texto(slide_cuerpo, titulo, 1.85, 0.4, font_size=18, color=[65,93,68], word_wrap= True, width_in=10, height_in=0.5, bold=True, align =  'center', font_name="Arial Black", shadow=True)
    agregar_texto(slide_cuerpo, titulo_2, 1.85, 0.65, font_size=16, color=[128, 0, 0], word_wrap= True, width_in=10, height_in=0.5, bold=True, align = 'center'  )

    
    
    agregar_texto(slide_cuerpo, fechas_final_final[2], 7.8, 1.4, font_size=16, color=[64, 64, 64], word_wrap= True, width_in=1.2, height_in=0.25, bold=True, align = 'center'  )
    agregar_texto(slide_cuerpo, fecha_inicial_anio, 9.1, 1.4, font_size=16, color=[64, 64, 64], word_wrap= True, width_in=1.2, height_in=0.25, bold=True, align = 'center'  )
       

    datos =datos_resultados[0]
    nun_x =2.083
    for x in datos:
        agregar_texto(slide_cuerpo, str(x), 7.8, nun_x, font_size=16, color=[64, 64, 64], word_wrap= True, width_in=1.2, height_in=0.25, bold=True, align = 'center'  )
        nun_x = nun_x + 0.257


    datos =datos_resultados[1]
    nun_x =2.083
    for x in datos:
        agregar_texto(slide_cuerpo, str(x), 9.1, nun_x, font_size=16, color=[64, 64, 64], word_wrap= True, width_in=1.2, height_in=0.25, bold=True, align = 'center'  )
        nun_x = nun_x + 0.257

    
    datos =datos_resultados[3]
    nun_x =2.083

    for x in datos:
        if int(x) >0:
            color=[0,176,80]            
        elif int(x) < 0:
            color=[125, 0, 0]
        else:
            color=[64, 64, 64]

        agregar_texto(slide_cuerpo, str(f"{float(x):.0f}"), 10.3, nun_x, font_size=16, color=color, word_wrap= True, width_in=1.2, height_in=0.25, bold=True, align = 'center'  )
        nun_x = nun_x + 0.257

    
    datos =datos_resultados[2]
    nun_x =2.083

   
    for x in datos:
        if int(x) >0:
            trinagulo_color = "tri_verde_new.png"
            
        elif int(x) < 0:
            trinagulo_color = "tri_rojo_new.png"
        else:
            trinagulo_color="tri_amarillo_new.png"
        tirangulo = f"{ruta_base}/static/img/{trinagulo_color}"

        agregar_texto(slide_cuerpo, str(f"{float(x):.0f} %"), 11.3, nun_x, font_size=16, color=[64, 64, 64], word_wrap= True, width_in=0.7, height_in=0.25, bold=True, align = 'right'  )

        slide_cuerpo.shapes.add_picture(
            tirangulo,
            Inches(12.2),   # left
            Inches(nun_x+0.05),   # top
            width=Inches(0.15)  # opcional
        )
        nun_x = nun_x + 0.257

    
    cocaina=datos_resultados[1][1]
    pbc=datos_resultados[1][3]
    marihuna=datos_resultados[1][2]

    agregar_texto(slide_cuerpo, str(cocaina), 7, 5.25, font_size=16, color=[255,255,255], word_wrap= True, width_in=1.5, height_in=0.3, bold=True, align = 'center', fill=True, fill_color=(56,87,35)  )
    agregar_texto(slide_cuerpo, str(pbc), 8.85, 5.25, font_size=16, color=[255,255,255], word_wrap= True, width_in=1.5, height_in=0.3, bold=True, align = 'center', fill=True, fill_color=(56,87,35)  )
    agregar_texto(slide_cuerpo, str(marihuna), 10.7, 5.25, font_size=16, color=[255,255,255], word_wrap= True, width_in=1.5, height_in=0.3, bold=True, align = 'center', fill=True, fill_color=(56,87,35)  )



    nota = f"Elaboración: propia DOCNA – {fecha_inicial_dia} de {fecha_inicial_mes} del {fecha_inicial_anio}"
    agregar_texto(slide_cuerpo, nota, 7.3, 7.3, font_size=9, color=[60,60,60], word_wrap= True, width_in=3, height_in=0.2, bold=False, align = '', fill=False, fill_color=(56,87,35)  )
 

    
    slide_portada = prs.slides.add_slide(prs.slide_layouts[7])
    link_2 = link #carpeta base
    nombre_carpeta = f"DOCNA/"
    link = os.path.join(link_2, nombre_carpeta)
    dirercion_archvios = link
    
    if os.path.exists(link_2):
        for elemento in os.listdir(link_2):
            ruta = os.path.join(link_2, elemento)
            if os.path.isdir(ruta):
                shutil.rmtree(ruta)
            else:
                os.remove(ruta)
    else:
        os.mkdir(link_2)
    

    if not os.path.exists(link):
        os.mkdir(link)

    # Guardar archivo
    direcion = link+str("test")+'.pptx'
    prs.save(direcion)

    direcion = link+str("test")+'.pptx'

    DIRECION_3 = os.getenv('DIRECION_3')
    direcion= DIRECION_3+nombre_carpeta+str("test")+'.pptx'
    nombre = f"2 INL {fecha_inicial_dia} - {fecha_inicial_mes}"
    return [direcion, nombre]