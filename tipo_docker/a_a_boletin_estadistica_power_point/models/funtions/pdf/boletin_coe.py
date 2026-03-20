from flask import make_response
from datetime import date, time, datetime
from tipo_docker.a_a_boletin_estadistica_power_point.models.estadistica.estadistica_boletin_coe import *
from __init__ import *
from tipo_docker.z_z_configuarcion.caligrafia import *
from tipo_docker.z_z_configuarcion.fechas import *
from tipo_docker.z_z_configuarcion.header import *
from tipo_docker.z_z_configuarcion.logo import *
from tipo_docker.z_z_configuarcion.titulos import *  
#funcion para estaditica

from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.enum.chart import XL_TICK_MARK
from pptx.util import Pt

from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.chart import XL_TICK_LABEL_POSITION

from pptx.enum.shapes import MSO_SHAPE

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_VERTICAL_ANCHOR
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.dml import MSO_COLOR_TYPE
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_FILL
from pptx.oxml.xmlchemy import OxmlElement



def cuadros_divisio(slide, valor_division, x, y):
    shapes = slide.shapes
    left = Inches(x)
    top = Inches(y)

    width = Inches(0.6)
    height = Inches(0.2)

    shape = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )

    # subtitle = slide.placeholders[1]
    text_frame = shape.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = str(valor_division)

    font = run.font
    font.name = 'Calibri'
    font.size = Pt(11)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_1


    # shape.text = text_frame
    # font = shape.font
    # p.alignment = PP_ALIGN.LEFT
    font.color.rgb = RGBColor(255, 255, 255)

    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(20, 62, 52)

def charts(slide, num_placeholder, categories, series_ad, series_nombre, afe):


    placeholder = slide.placeholders[num_placeholder]
    chart_data = ChartData()
    chart_data.categories = categories
    chart_data.add_series(series_nombre,series_ad)
    
    
    graphic_frame = placeholder.insert_chart(XL_CHART_TYPE.DOUGHNUT, chart_data)
    chart = graphic_frame.chart
    chart.chart_type

    chart.has_legend = True
    chart.font.size = Pt(14)
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.font.color.rgb = RGBColor(70, 70, 70)



    categorias = chart.plots[0]
    categorias.has_text_frame = False

    plot = chart.plots[0]
    plot.has_data_labels = True

    data_labels = plot.data_labels

    data_labels.font.size = Pt(14)
    data_labels.font.color.rgb = RGBColor(125, 0, 0)
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    
    chart.has_title = False

    colors = {
        series_nombre: ['3fb872', '035827', '0d2116'],  # Example colors
        # Add more series names and color lists as needed
    }
    for series in chart.series:
        if series.name in colors:
            for i, color_code in enumerate(colors[series.name]):
                if i < len(series.points):
                    fill = series.points[i].format.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor.from_string(color_code)
        i = 0
        series = chart.series[0]
        for point in series.points:
            point.data_label.has_text_frame = True
            #Assigning custom text for data label associated with each data-point
            point.data_label.text_frame.text =  str(int(chart.series[0].values[i])) +" "+afe
       
            for run in point.data_label.text_frame.paragraphs[0].runs:
                run.font.size = Pt(12)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0, 0, 0)
            i+=1
    
def tabla(slide, datos, anio_ant, anio_act, num_placeholder, afectacion):
    
    tamnio = len(datos)

    placeholder = slide.placeholders[num_placeholder]
    graphic_frame = placeholder.insert_table(rows=tamnio+2, cols=6)
    table = graphic_frame.table

    start_cell = table.cell(0, 0)
    end_cell = table.cell(0,5)
    start_cell.merge(end_cell)

    cell = table.cell(0, 0)
    cell.text = afectacion
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.size = Pt(12)
    paragraph.font.bold = True

    cell = table.cell(1, 0)
    cell.text = 'DIVISIÓN'
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.size = Pt(12)
    paragraph.font.bold = True
  
    cell = table.cell(1, 1)
    cell.text = str(anio_ant)
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.size = Pt(12)
    paragraph.font.bold = True

    cell = table.cell(1, 2)
    cell.text = str(anio_act)
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.size = Pt(12)
    paragraph.font.bold = True

    cell = table.cell(1, 3)
    cell.text = str("DIF.")
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.size = Pt(12)
    paragraph.font.bold = True
        
    cell = table.cell(1, 4)
    cell.text = str("%")
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.size = Pt(12)
    paragraph.font.bold = True
            
    cell = table.cell(1, 5)
    cell.text = str("SEM.")
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.size = Pt(12)
    paragraph.font.bold = True

    table.columns[0].width = Inches(1)
    table.columns[1].width = Inches(0.5)
    table.columns[2].width = Inches(0.6)
    table.columns[3].width = Inches(0.6)
    table.columns[4].width = Inches(0.6)
    table.columns[5].width = Inches(0.6)


    numero = 1
    fila = 2
    
    tamnio = len(datos)
    for x in datos:
            
            if fila <= tamnio+1:
                row = table.rows[fila]
                row.height = Inches(0.1)

                cols = table.columns[1]
                cols.width = Inches(1)

                cell = table.cell(fila, 0)
                cell.text = x[0]
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(9)
                paragraph.font.bold = True

                cell = table.cell(fila, 1)
                cell.text = str(x[1])
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(9)
                paragraph.font.bold = True
                
                cell = table.cell(fila, 2)
                cell.text = str(x[2])
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(9)
                paragraph.font.bold = True
                                
                cell = table.cell(fila, 3)
                cell.text = str(x[3])
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(9)
                paragraph.font.bold = True
                
                cell = table.cell(fila, 4)
                cell.text = str(x[4])
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(9)
                paragraph.font.bold = True
                                
                cell = table.cell(fila, 5)
                cell.text = str(x[5])
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(9)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(255, 255, 255) 

                fila= fila+1

    r=2
    for x in datos:
        c=5
        for y in x:
            
            #Insert data into table
            if x[5] == "-":
                cell = table.cell(r, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(128, 0, 0)    
            elif x[5] == "+":
                cell = table.cell(r, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(55, 86, 35)   
            else:
                cell = table.cell(r, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(255, 192, 0)     
            cell = apply_border(cell, edges = ["left", "right", "top", "bottom"])

        r = r +1

def fucion_para_crear_tabla(slide, data,left, top, width, height, numero ):

    shape = slide.shapes.add_table(rows=len(data), cols=len(data[0]), left=left, top=top, width=width, height=height)
    table = shape.table

    # Llenar la tabla con datos
    
    for row_index, row_data in enumerate(data):
        for col_index, cell_data in enumerate(row_data):

            table.cell(row_index, col_index).text = cell_data
            table.cell(row_index, col_index).vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            table.cell(row_index, col_index).text_frame.paragraphs[0].font.size =Pt(8)
            table.cell(row_index, col_index).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            table.cell(row_index, 0).text_frame.paragraphs[0].font.bold = True
            table.cell(row_index, numero).text_frame.paragraphs[0].font.bold = True
            table.cell(row_index, col_index).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
            table.rows[row_index].height = Inches(0.1) 
            table.columns[col_index].width = Inches(0.58)
            
            table.cell(row_index, numero).fill.solid()
            table.cell(row_index, numero).fill.fore_color.rgb = pptx.dml.color.RGBColor(220, 220, 220)  

            table.cell(0, col_index).fill.solid()
            table.cell(0, col_index).fill.fore_color.rgb = pptx.dml.color.RGBColor(128, 0, 0)
            table.cell(0, col_index).text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)  
            table.cell(row_index, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    table.columns[0].width = Inches(1.5)  

def estadistica_resaltantes(fecha_inicial_u_l, fecha_final_u_l, filtro, dirercion_archvios, nombre_carpeta):

    
    #-----------------------------------
    #modulo para el calculo de los resultados 
    fechas_inicial = fecha(fecha_inicial_u_l)
    fechas_final = fecha(fecha_final_u_l)
    
    fecha_titulo = fechas_inicial[0]+" " +  fechas_inicial[1] + " AL " + fechas_final[0] + " "+fechas_final[1] + " "+fechas_final[2]
    fecha_titulo = fecha_titulo.upper()

    fecha_dia = fechas_final[0] + " - "+fechas_final[4] + " - "+fechas_final[2]

    fecha_dia= fecha_dia.upper()

    fecha_inicial_dia_inicial = fechas_inicial[0]
    fecha_inicial_mes_inicial = fechas_inicial[1]
    fecha_inicial_anio_inicial = fechas_inicial[2]

    fecha_inicial_dia = fechas_final[0]
    fecha_inicial_mes = fechas_final[1]
    fecha_inicial_anio = fechas_final[2]

    calcular_spoa = Calculo_Spoa(fecha_inicial_u_l, fecha_final_u_l, filtro)
    RESULTADOS = calcular_spoa.calculo_boletin_cuadro()

    prs = Presentation('C:/Users/nicolas.valbuena/Documents/programacion 2023/server/tipo_docker/a_a_boletin_estadistica_power_point/models/documentos/template/templete.pptx')
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))

    titulo = "EVALUACIÓN OBJETIVO ESTRATÉGICO #2 DEBILITAR LAS CAPACIDADES DE LA AMENAZA POR DIVISIONES."
    placeholder = slide.placeholders[0]
    placeholder = slide.shapes.title
    placeholder.text = titulo

    placeholder = slide.placeholders[10]
    text_frame = placeholder.text_frame
    text_frame.text = fecha_dia
    slide = prs.slides.add_slide(prs.slide_layouts[2])

    #left =  Inches(2)
    #top = Inches(6)
    #widt = Inches(1)
    #height = Inches(1)
    #try:
       # slide.shapes.add_picture("C:/Users/nicolas.valbuena/Documents/programacion 2023/server/static/img/resultados/asesinados.png", left, top, widt, height)
    #except FileNotFoundError:
       # print("La imagen no se encontro en la ruta especificada")
    #except Exception as e:
       # print("ocurrio un error al insertar la imagen :{e}")
    

    
    titulo = "EVALUACIÓN OBJETIVO ESTRATÉGICO #2 DEBILITAR LAS CAPACIDADES DE LA AMENAZA POR DIVISIONES."
    placeholder = slide.placeholders[0]
    placeholder = slide.shapes.title
    placeholder.text = titulo

    placeholder = slide.placeholders[21]
    text_frame = placeholder.text_frame
    text_frame.text = fecha_titulo

    # Definir la posición y tamaño de la tabla
    top = Inches(1.5)
    left = Inches(0.1)
    width = Inches(12)
    height = Inches(3)

    # Datos para la tabla
    data = []
    data_1 = []
    data_2 = []
    data_3 = []
    data_4 = []
    numero=1
    data.append(RESULTADOS[0])
    data_1.append(RESULTADOS[0])
    data_2.append(RESULTADOS[0])
    data_3.append(RESULTADOS[0])
    data_4.append(RESULTADOS[0])
    
    columna=0
    for x in RESULTADOS[0]:
        if x == "TOTAL":
            columna = columna
            break
        columna = columna +1

    for x in RESULTADOS:

        if numero <=17:
            data.append(RESULTADOS[numero])
        elif numero >=18 and numero <=35:
            data_1.append(RESULTADOS[numero])
        elif numero >=36 and numero <=55:
            data_2.append(RESULTADOS[numero])
        elif numero >=56 and numero <=75:
            data_3.append(RESULTADOS[numero])
        elif numero >=76 and numero <=91:
            data_4.append(RESULTADOS[numero])
        numero=  numero+1
   
    fucion_para_crear_tabla(slide, data,left, top, width, height, columna)

    slide = prs.slides.add_slide(prs.slide_layouts[2])
    placeholder = slide.placeholders[0]
    placeholder = slide.shapes.title
    placeholder.text = titulo

    placeholder = slide.placeholders[21]
    text_frame = placeholder.text_frame
    text_frame.text = fecha_titulo

    # Definir la posición y tamaño de la tabla

    # Datos para la tabla
    # Datos para la tabla

    fucion_para_crear_tabla(slide, data_1,left, top, width, height, columna )

    slide = prs.slides.add_slide(prs.slide_layouts[2])
    placeholder = slide.placeholders[0]
    placeholder = slide.shapes.title
    placeholder.text = titulo

    placeholder = slide.placeholders[21]
    text_frame = placeholder.text_frame
    text_frame.text = fecha_titulo

    fucion_para_crear_tabla(slide, data_2,left, top, width, height, columna )

    slide = prs.slides.add_slide(prs.slide_layouts[2])
    placeholder = slide.placeholders[0]
    placeholder = slide.shapes.title
    placeholder.text = titulo

    placeholder = slide.placeholders[21]
    text_frame = placeholder.text_frame
    text_frame.text = fecha_titulo

    fucion_para_crear_tabla(slide, data_3,left, top, width, height, columna )


    slide = prs.slides.add_slide(prs.slide_layouts[2])
    placeholder = slide.placeholders[0]
    placeholder = slide.shapes.title
    placeholder.text = titulo

    placeholder = slide.placeholders[21]
    text_frame = placeholder.text_frame
    text_frame.text = fecha_titulo

    fucion_para_crear_tabla(slide, data_4,left, top, width, height, columna )

    direcion = dirercion_archvios+str("test")+'.pptx'
    prs.save(direcion)

    direcion = dirercion_archvios+str("test")+'.pptx'

    DIRECION_3 = os.getenv('DIRECION_3')
    direcion= DIRECION_3+nombre_carpeta+str("test")+'.pptx'
    
    return [direcion, "nombre_doc"]
