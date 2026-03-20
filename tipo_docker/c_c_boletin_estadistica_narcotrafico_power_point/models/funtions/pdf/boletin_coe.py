from flask import make_response
from datetime import date, time, datetime

from tipo_docker.c_c_boletin_estadistica_narcotrafico_power_point.models.estadistica.estadistica_boletin_coe import *

from __init__ import *

from tipo_docker.z_z_configuarcion.caligrafia import *
from tipo_docker.z_z_configuarcion.fechas import *
from tipo_docker.z_z_configuarcion.header import *
from tipo_docker.z_z_configuarcion.logo import *
from tipo_docker.z_z_configuarcion.titulos import *  
from tipo_docker.z_z_configuarcion.tipo_ayudas import *

from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.enum.chart import XL_TICK_MARK
from pptx.util import Pt

from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.chart import XL_TICK_LABEL_POSITION

from pptx.enum.shapes import MSO_SHAPE

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_VERTICAL_ANCHOR
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.dml import MSO_COLOR_TYPE
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_FILL
from pptx.oxml.xmlchemy import OxmlElement
import os
import shutil



def set_cell_border(cell):
    """Aplica bordes negros a una celda de tabla en pptx."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for side in ['L', 'R', 'T', 'B']:
        ln = OxmlElement(f'a:ln{side}')
        ln.set('w', '12700')  # 1pt en EMUs

        solidFill = OxmlElement('a:solidFill')
        srgbClr = OxmlElement('a:srgbClr')
        srgbClr.set('val', '000000')
        solidFill.append(srgbClr)
        ln.append(solidFill)

        prstDash = OxmlElement('a:prstDash')
        prstDash.set('val', 'solid')
        ln.append(prstDash)

        tcPr.append(ln)




def crear_tabla_estilizada(prs, slide, data, size_pt, left=1, top=1, width=8):
    """
    Inserta una tabla con estilo en una diapositiva con filas de 0.5 cm de alto (≈0.19685 pulgadas),
    sin espacios verticales ni márgenes en el texto.
    """
    rows, cols = len(data), len(data[0])
    row_height = Inches(0.19685)  # 0.5 cm
    total_height = row_height * rows

    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), total_height)
    table = table_shape.table

    # Ancho de columnas (ajustable)
    col_widths = [Inches(1.7), Inches(0.7), Inches(0.7), Inches(0.6)]
    for i, w in enumerate(col_widths[:cols]):
        table.columns[i].width = w

    # Altura de filas
    for i in range(rows):
        table.rows[i].height = row_height

    # Colores
    header_bg = RGBColor(128, 0, 0)
    white = RGBColor(255, 255, 255)
    black = RGBColor(0, 0, 0)

    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val.strip()

            # 🔹 Eliminar márgenes internos
            cell.text_frame.margin_top = 0
            cell.text_frame.margin_bottom = 0
            cell.text_frame.margin_left = 0
            cell.text_frame.margin_right = 0

            # 🔹 Anclar verticalmente al centro
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 🔹 Quitar espaciado antes/después del párrafo
            para = cell.text_frame.paragraphs[0]
            para.space_before = Pt(0)
            para.space_after = Pt(0)
            para.line_spacing = 1  # línea ajustada, sin extra

            # 🔹 Alineación horizontal
            if i == 0:
                para.alignment = PP_ALIGN.CENTER
            elif j == 0:
                para.alignment = PP_ALIGN.LEFT
            elif j == len(row) - 1:
                para.alignment = PP_ALIGN.RIGHT
            else:
                para.alignment = PP_ALIGN.CENTER

            # 🔹 Fuente
            run = para.runs[0] if para.runs else para.add_run()
            run.font.name = "Arial"
            run.font.size = Pt(size_pt)
            run.font.bold = (i == 0)
            run.font.color.rgb = white if i == 0 else black

            # 🔹 Fondo
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = header_bg if i == 0 else white

            # 🔹 Bordes
            apply_border(cell)

    return table_shape

def crear_tabla_estilizada_con_escudos(prs, slide, data, size_pt, ruta_base, escudos, left=1, top=1, width=8):
    """
    Inserta una tabla con estilo y coloca imágenes (círculos) en la última columna por fila.
    """
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor

    row_height = Inches(0.19685)  # 0.5 cm
    rows, cols = len(data), len(data[0])
    total_height = row_height * rows

    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), total_height)
    table = table_shape.table

    # Definir anchos fijos por columna
    col_widths = [Inches(1.45), Inches(0.65), Inches(0.65), Inches(0.6)]
    for i, w in enumerate(col_widths[:cols]):
        table.columns[i].width = w

    for i in range(rows):
        table.rows[i].height = row_height

    # Colores
    header_bg = RGBColor(128, 0, 0)
    white = RGBColor(255, 255, 255)
    black = RGBColor(0, 0, 0)

    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val

            cell.text_frame.margin_top = 0
            cell.text_frame.margin_bottom = 0
            cell.text_frame.margin_left = 0
            cell.text_frame.margin_right = 0
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            para = cell.text_frame.paragraphs[0]
            para.space_before = Pt(0)
            para.space_after = Pt(0)

            if i == 0:
                para.alignment = PP_ALIGN.CENTER
            elif j == 0:
                para.alignment = PP_ALIGN.LEFT
            elif j == len(row) - 1:
                para.alignment = PP_ALIGN.RIGHT
            else:
                para.alignment = PP_ALIGN.CENTER

            run = para.runs[0] if para.runs else para.add_run()
            run.font.name = "Arial"
            run.font.size = Pt(size_pt)
            run.font.bold = (i == 0)
            run.font.color.rgb = white if i == 0 else black

            # Colores de fondo
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = header_bg if i == 0 else white

            # Bordes
            apply_border(cell, edges=["left", "right", "top", "bottom"])

        # Insertar imagen en la última columna (excepto encabezado)
        row_height = Inches(0.235)  # 0.5 cm
        if i > 0 and i - 1 < len(escudos):
            try:
                img_path = f"{ruta_base}/static/img/escudos/{escudos[i - 1]}.png"
                img_width = Inches(0.18)
                img_height = Inches(0.18)

                # Posición base de tabla en pulgadas
                table_left = Inches(left)
                table_top = Inches(top)

                # Posición izquierda de la última celda
                cell_left = table_left + sum(col_widths[:cols - 1])
                cell_top = table_top + row_height * i

                desplazar_derecha = Inches(0.2)  # mover a la derecha (≈0.8 mm)
                desplazar_abajo = Inches(0.015)    # mover hacia abajo (≈0.5 mm)

                # Centrado de la imagen en la celda
                img_left = cell_left + col_widths[cols - 1] / 2 - img_width / 3 - desplazar_derecha
                img_top = cell_top + row_height / 2 - img_height / 2 + desplazar_abajo

                slide.shapes.add_picture(img_path, img_left, img_top, width=img_width, height=img_height)

            except Exception as e:
                print(f"Error insertando imagen: {img_path} -> {e}")

    return table_shape

def crear_tabla_estilizada_con_escudos_5(prs, slide, data, size_pt, ruta_base, escudos, left=1, top=1, width=8):
    """
    Inserta una tabla con estilo (5 columnas) y coloca imágenes (círculos) en la última columna por fila (excepto encabezado).
    """
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor

    row_height = Inches(0.25)  # 0.5 cm
    rows, cols = len(data), 5  # Forzar 5 columnas
    total_height = row_height * rows

    # Crear tabla con 5 columnas
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), total_height)
    table = table_shape.table

    # Anchos de columnas: 4 para datos, 1 para imagen
    col_widths = [Inches(1.42), Inches(0.64), Inches(0.64), Inches(0.66), Inches(0.55)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for i in range(rows):
        table.rows[i].height = row_height

    # Colores
    header_bg = RGBColor(128, 0, 0)
    white = RGBColor(255, 255, 255)
    black = RGBColor(0, 0, 0)

    for i, row in enumerate(data):
        for j in range(cols):
            cell = table.cell(i, j)
            
            if j < len(row):
                cell.text = row[j]
            else:
                cell.text = ""  # Celda vacía en columna 5 (imagen)

            # Estilo y márgenes
            cell.text_frame.margin_top = 0
            cell.text_frame.margin_bottom = 0
            cell.text_frame.margin_left = 0
            cell.text_frame.margin_right = 0
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            para = cell.text_frame.paragraphs[0]
            para.space_before = Pt(0)
            para.space_after = Pt(0)

            if i == 0:
                para.alignment = PP_ALIGN.CENTER
            elif j == 0:
                para.alignment = PP_ALIGN.LEFT
            elif j == cols - 1:
                para.alignment = PP_ALIGN.RIGHT
            else:
                para.alignment = PP_ALIGN.CENTER

            run = para.runs[0] if para.runs else para.add_run()
            run.font.name = "Calibri"
            run.font.size = Pt(size_pt)
            run.font.bold = (i == 0)
            run.font.color.rgb = white if i == 0 else black

            # Fondo
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = header_bg if i == 0 else white

            # Bordes
            apply_border(cell, edges=["left", "right", "top", "bottom"])

        # Insertar imagen en última columna (solo si hay escudo y no es encabezado)

        texto = table.cell(0, 2).text
        ancho_col = col_widths[2]  # columna 3



        row_height = Inches(0.25)  # 0.5 
        if i > 0 and i - 1 < len(escudos):
            try:
                img_path = f"{ruta_base}/static/img/escudos/{escudos[i - 1]}.png"
                img_width = Inches(0.18)
                img_height = Inches(0.18)

                # Posición base
                table_left = Inches(left)
                table_top = Inches(top)

                # Posición de la última columna
                cell_left = table_left + sum(col_widths[:4])
                cell_top = table_top + row_height * i

                # Ajustes finos
                desplazar_derecha = Inches(0.2)
                
     
                if len(texto) > 10:
                    desplazar_abajo = Inches(0.32)
                else:
                    desplazar_abajo = Inches(0.16)

                img_left = cell_left + col_widths[4] / 2 - img_width / 2 - desplazar_derecha
                img_top = cell_top + row_height / 2 - img_height / 2 + desplazar_abajo

                slide.shapes.add_picture(img_path, img_left, img_top, width=img_width, height=img_height)

            except Exception as e:
                print(f"Error insertando imagen: {img_path} -> {e}")

    return table_shape


def apply_border(cell, edges=["left", "right", "top", "bottom"], border_color="000000", border_width=1):
    if type(edges) is not list:
        edges = [edges]
    border_width = str(border_width * Pt(1))

    def SubElement(parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

    lines = [{"left": 'a:lnL',
              "right": 'a:lnR',
              "top": 'a:lnT',
              "bottom": 'a:lnB'}[_] for _ in edges]

    if hasattr(cell.fill, "fore_color"):
        fill_color = cell.fill.fore_color.rgb
    cell.fill.background()

    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for line in lines:
        # Remove duplicate tag if it exists
        tag = line.split(":")[-1]
        for e in tcPr.getchildren():
            if tag in str(e.tag):
                tcPr.remove(e)

        ln = SubElement(tcPr, line, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(ln, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        SubElement(ln, 'a:prstDash', val='solid')
        SubElement(ln, 'a:round')
        SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')

    cell.fill.solid()
    if cell.fill.type == MSO_FILL.SOLID:
        cell.fill.fore_color.rgb = fill_color
    return cell

def cuadros_divisio(slide, valor_division, x, y):
  # ... dentro de tu función o flujo
    shapes = slide.shapes
    left = Inches(x)
    top = Inches(y)

    width = Inches(1.1)
    height = Inches(0.3)

    shape = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )

    # Eliminar fondo del rectángulo
    shape.fill.background()

    # Eliminar bordes
    shape.line.fill.background()

    # Agregar texto
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = str(valor_division)

    # Estilo de texto
    font = run.font
    font.name = 'Arial'
    font.size = Pt(14)
    font.bold = True
    font.color.rgb = RGBColor(0, 0, 0)

    # Alineación opcional
    p.alignment = PP_ALIGN.LEFT

def cuadros_divisio_r(slide, valor_division, x, y, rotar, color, ancho):
    shapes = slide.shapes
    left = Inches(x)
    top = Inches(y)
    width = Inches(ancho)
    height = Inches(0.3)

    shape = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )

    # Eliminar fondo y bordes
    shape.fill.background()
    shape.line.fill.background()

    # Agregar texto
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = str(valor_division)

    # Estilo de texto
    font = run.font
    font.name = color[3]
    font.size = Pt(color[4])
    font.bold = color[5]
    font.italic = color[6]
    font.color.rgb = RGBColor(color[0],color[1],color[2])

    # Alineación
    p.alignment = PP_ALIGN.CENTER

    # 🔄 Rotar el texto o el shape (90 grados)
    shape.rotation = rotar  # rota todo el cuadro de texto

    # Si quieres solo el texto en vertical sin girar la forma:
    # text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    # text_frame.text_anchor = MSO_VERTICAL_ANCHOR.VERTICAL


# Función para agregar línea personalizada
def agregar_linea(slide, left_in, top_in, largo_in, orientacion="horizontal", color_rgb=(0, 0, 0), grosor_pt=2):
    left = Inches(left_in)
    top = Inches(top_in)

    if orientacion == "horizontal":
        width = Inches(largo_in)
        height = Inches(0)  # Línea horizontal
    elif orientacion == "vertical":
        width = Inches(0)
        height = Inches(largo_in)  # Línea vertical
    elif orientacion == "diagonal":
        width = Inches(largo_in)
        height = Inches(largo_in)  # Diagonal con 45°
    else:
        raise ValueError("Orientación no válida. Usa: horizontal, vertical o diagonal.")

    line = slide.shapes.add_shape(
        autoshape_type_id=1,  # Línea
        left=left,
        top=top,
        width=width,
        height=height
    )

    # Estilizar línea
    line.line.color.rgb = RGBColor(*color_rgb)
    line.line.width = Pt(grosor_pt)


from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def agregar_texto(
    slide,
    texto,
    left_in,
    top_in,
    width_in=0.5,
    height_in=0.3,
    font_size=12,
    font_name="Calibri",
    color=(127, 127, 127),
    bold=False,
    italic=False,
    align="left",
    word_wrap=False
):

    left = Inches(left_in)
    top = Inches(top_in)
    width = Inches(width_in)
    height = Inches(height_in)

    textbox = slide.shapes.add_textbox(left, top, width, height)

    tf = textbox.text_frame

    # 🔥 Configuración correcta
    tf.word_wrap = word_wrap
    tf.auto_size = None
    tf.clear()

    # 🔥 Quitar márgenes internos
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)

    # 🔥 Usar primer párrafo (NO add_paragraph)
    p = tf.paragraphs[0]
    p.text = texto

    # Alineación
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    elif align == "justify":
        p.alignment = PP_ALIGN.JUSTIFY
    else:
        p.alignment = PP_ALIGN.LEFT

    # Fuente
    font = p.font
    font.size = Pt(font_size)
    font.name = font_name
    font.bold = bold
    font.italic = italic
    font.color.rgb = RGBColor(*color)

def pdf_boletin_narcotrafico(fecha_inicial_u_l, fecha_final_u_l, filtro, link, puerto, fecha_primer_lapso_final, fecha_ultimo_lapso_final):
    fecha_dt = datetime.strptime(fecha_final_u_l, '%Y-%m-%d')
    fecha_inicio = fecha(fecha_inicial_u_l)
    fecha_final = fecha(fecha_final_u_l)
    
    fecha_palso_anterior_i = fecha(fecha_primer_lapso_final)
    fecha_palso_anterior_f = fecha(fecha_final_u_l)

    anio_inicio_i_f = fecha_palso_anterior_i[2]
    anio_inicio_f_f = fecha_palso_anterior_f[2]

    dia_inicio_f, mes_inicio_f, anio_inicio_f = fecha_palso_anterior_i[0], fecha_palso_anterior_i[1], fecha_palso_anterior_i[2]

    dia_inicio, mes_inicio, anio_inicio = fecha_final[0], fecha_final[1], fecha_final[2]
    dia_fin, mes_fin, anio_fin, mes_fin_numero = fecha_final[0], fecha_final[1], fecha_final[2], fecha_final[3]
    mes_num_inicio = fecha_final[1]
    dia_sem = dias_semana(fecha_dt.weekday())
    fecha_dia_anterior = f"{anio_inicio}-01-01"

    calcular_spoa = Calculo_Spoa(fecha_inicial_u_l, fecha_final_u_l, fecha_dia_anterior, filtro, str(anio_inicio), str(mes_fin.upper()), int(mes_fin_numero), fecha_primer_lapso_final, fecha_ultimo_lapso_final)
    resultado = calcular_spoa.resultados_narcotrafico_valores()
    

    # Crear presentación y título
    #ruta del template de la ayuda de narcotrafico img_path = f"{ruta_base}/static/img/escudos/{escudos[i - 1]}.png"
    ruta_base = filtro[15]
    ayuda = f"{ruta_base}/static/template/template_ayuda_docna_2026.pptx"
    prs = Presentation(ayuda)
    #este codigo es para saber que tipo de ayudas tengo
    #tipo_ayudas_template(prs)
   
    # Slide de portada
    #slide_portada = prs.slides.add_slide(prs.slide_layouts[0])
    #slide_portada = prs.slides.add_slide(prs.slide_layouts[1])
    slide_portada = prs.slides.add_slide(prs.slide_layouts[2])
    #titulo_slide = slide_portada.shapes.title
    #titulo_texto = f"EVALUACIÓN LÍNEA DE OPERACIONES “PAYA” \n \n{dia_inicio} - {mes_num_inicio} - {anio_inicio}"
    #titulo_slide.text = titulo_texto.upper()

    # Slide de resultados
    
    slide_result = prs.slides.add_slide(prs.slide_layouts[3])
    #slide_result.shapes.title.text = f"EVALUACIÓN LÍNEA DE OPERACIONES “PAYA” ".upper()
        # colocar fecha en la ayuda
    fecha_dia_anterior = (str((int(dia_inicio)-1)) +" - "+ str(mes_num_inicio)  +" - "+ str(anio_inicio)).upper()
    fecha_actual = (f"{(int(dia_inicio_f))} de {mes_inicio_f} hasta el  {(int(dia_inicio))} - {mes_num_inicio} - {anio_inicio}").upper()

    agregar_texto(slide_result, fecha_actual, 7.5, 0.88)
    # Cuadros resumen
    posiciones_y = [1.65, 2.6, 3.55, 4.5, 5.45, 6.03, 6.4]
    for i, y in enumerate(posiciones_y):
        cuadros_divisio(slide_result, resultado[0][i], 1.8, y)

    left = Inches(3.55)     # distancia desde la izquierda
    top = Inches(-0.1)      # distancia desde arriba
    width = Inches(6.81)    # ancho de la imagen (opcional)
    height = Inches(8.47)   # alto de la imagen (opcional)

    mapa_narcotraficos = '{}static/img/img_mapas/mapa_narcotrafico.png'.format(filtro[15])
    slide_result.shapes.add_picture(mapa_narcotraficos, left, top, width=width, height=height)
    # Tabla de resultados
    crear_tabla_estilizada(prs, slide_result, resultado[1], 10, left=9.35, top=5.5, width=3)
    # Coordenadas base
    left = Inches(12.47)
    top_base = 5.8
    step = 0.28  # diferencia vertical entre imágenes
    width = Inches(0.18)
    height = Inches(0.18)

    # Insertar cada escudo
    for i in range(len(resultado[2])):
        top = Inches(top_base + i * step)
        escudo_path = f'{filtro[15]}static/img/escudos/{resultado[2][i]}.png'
        slide_result.shapes.add_picture(escudo_path, left, top, width=width, height=height)

    # Slide lamina de cuadros 
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    #slide.shapes.title.text = f"EVALUACIÓN LÍNEA DE OPERACIONES “PAYA” ".upper()
    titulo = f"""EVALUACIÓN LÍNEA DE OPERACIONES “PAYA” CONSOLIDADO ANUAL""" 
    agregar_texto(slide, titulo, 1.8, 0.2, font_size=28, color=[56, 87, 35], word_wrap=True, width_in=10, bold=True)

    agregar_texto(slide, fecha_actual, 8, 0.88)
    resultado = calcular_spoa.resultados_narcotrafico_boletin()
    

    # Coordenadas de bloques de tabla
    bloques = [
        (0, 5,    0.9, 1.2),   # Fila 1 - izquierda
        (6, 11,   5.05, 1.2),  # Fila 1 - centro
        (12, 17,  9.2, 1.2),   # Fila 1 - derecha

        (18, 23,  0.9, 2.7),   # Fila 2 - izquierda
        (24, 29,  5.05, 2.7),  # Fila 2 - centro
        (30, 35,  9.2, 2.7),   # Fila 2 - derecha

        (36, 41,  0.9, 4.2),   # Fila 3 - izquierda
        (42, 47,  5.05, 4.2),  # Fila 3 - centro
        (48, 53,  9.2, 4.2),   # Fila 3 - derecha

        (54, 59,  3, 5.7),   # Fila 4 - izquierda
        (60, 65,  7.15, 5.7),  # Fila 4 - centro
        # (66, 71, 9.2, 5.7),  # Fila 4 - derecha (omitido)
    ]


    # Para cada bloque, pasa también las rutas de escudos
    for i, (inicio, fin, left, top) in enumerate(bloques):
        data = resultado[0][inicio:fin + 1]
        escudos_tabla = resultado[1][i * 5 : (i + 1) * 5]
        crear_tabla_estilizada_con_escudos(prs, slide, data, 8.5, filtro[15], escudos_tabla, left=left, top=top, width=4)

    escudos_img = [
        "escudo_ejc.jpg", "escudo_davaa.jpg", "escudo_div01.jpg",
        "escudo_div02.jpg", "escudo_div03.jpg", "escudo_div04.jpg",
        "escudo_div05.jpg", "escudo_div06.jpg", "escudo_div07.jpg",
        "escudo_div08.jpg", "escudo_futco.jpg"
    ]

    # Coordenadas de cada escudo (fila, columna)
    # Cada tupla representa: (fila, columna)
    posiciones = [
        (0, 0), (0, 1), (0, 2),
        (1, 0), (1, 1), (1, 2),
        (2, 0), (2, 1), (2, 2),
        (3, 0), (3, 1)
    ]

    # Base de coordenadas
    top_base = Inches(1.19)
    left_base = Inches(0.38)
    col_step = Inches(4.16)   # separación horizontal
    row_step = Inches(1.5)    # separación vertical
    width = Inches(0.5)
    height = Inches(0.5)

    # Insertar escudos
    for idx, (fila, columna) in enumerate(posiciones):
        if idx >= len(escudos_img):
            break

        # Posición base normal
        top = top_base + row_step * fila
        left = left_base + col_step * columna

        # 👉 Si es la última fila (fila 3), desplazar a la derecha
        if fila == 3:
            left += Inches(2.1)

        escudo_path = f'{filtro[15]}static/img/escudos/{escudos_img[idx]}'
        try:
            slide.shapes.add_picture(
                escudo_path,
                left,
                top,
                width=width,
                height=height
            )
        except Exception as e:
            print(f"Error insertando {escudo_path}: {e}")


    # Slide lamina de cuadros 
    slide = prs.slides.add_slide(prs.slide_layouts[4])


    

    #slide.shapes.title.text = f"EVALUACIÓN LÍNEA DE OPERACIONES “PAYA” ".upper()
    titulo = f"""EVALUACIÓN LÍNEA DE OPERACIONES “PAYA” META MENSUAL "{mes_num_inicio}" """.upper() 
    agregar_texto(slide, titulo, 2, 0.2, font_size=28, color=[56, 87, 35], word_wrap= True, width_in=10, height_in=1, bold=True)
    agregar_texto(slide, fecha_actual, 8, 0.88)

    resultado = calcular_spoa.resultados_narcotrafico_boletin_metas()
    
        # Coordenadas de bloques de tabla
    bloques = [
        (0, 5,    0.7, 2.1),   # Fila 1 - izquierda
        (6, 11,   4.7, 2.1),  # Fila 1 - centro
        (12, 17,  8.75, 2.1),   # Fila 1 - derecha

        (18, 23,  0.7, 4.7),   # Fila 2 - izquierda
        (24, 29,  4.7, 4.7),  # Fila 2 - centro
        (30, 35,  8.75, 4.7),   # Fila 2 - derecha

        #(36, 41,  0.9, 4.2),   # Fila 3 - izquierda
        #(42, 47,  5.05, 4.2),  # Fila 3 - centro
        #(48, 53,  9.2, 4.2),   # Fila 3 - derecha

        #(54, 59,  0.9, 5.7),   # Fila 4 - izquierda
        #(60, 65,  5.05, 5.7),  # Fila 4 - centro
        # (66, 71, 9.2, 5.7),  # Fila 4 - derecha (omitido)
    ]


    # Para cada bloque, pasa también las rutas de escudos
    for i, (inicio, fin, left, top) in enumerate(bloques):
        data = resultado[0][inicio:fin + 1]
        escudos_tabla = resultado[1][i * 5 : (i + 1) * 5]
        crear_tabla_estilizada_con_escudos_5(prs, slide, data, 8, filtro[15], escudos_tabla, left=left, top=top, width=4)
    # Ejemplo de uso:
    

    # === Lista de escudos ===
    escudos_img = [
        "ejc.png", "davaa_ayuda.png", "div01_ayuda.png",
        "div02_ayuda.png", "div03_ayuda.png", "div04_ayuda.png"
    ]

    # === Configuración general ===
    largo_linea = 1.5
    distancia_grupo = 2.35
    base_left = 0.7
    top_lineas = [1.7, 4.4]
    top_imagenes = [1.4, 4.1]
    img_width = Inches(0.5)
    img_height = Inches(0.5)

    # === Inserción automática ===
    indice_escudo = 0
    for fila in range(2):  # dos filas de líneas e imágenes
        top_linea = top_lineas[fila]
        top_img = Inches(top_imagenes[fila])
        left = base_left

        for _ in range(3):  # 3 bloques por fila
            # Línea izquierda
            agregar_linea(slide, left_in=left, top_in=top_linea, largo_in=largo_linea,
                        orientacion="horizontal", color_rgb=(51, 63, 80), grosor_pt=3)

            # Línea derecha
            right_line = left + distancia_grupo
            agregar_linea(slide, left_in=right_line, top_in=top_linea, largo_in=largo_linea,
                        orientacion="horizontal", color_rgb=(51, 63, 80), grosor_pt=3)

            # Imagen centrada entre líneas
            if indice_escudo < len(escudos_img):
                escudo_path = f'{filtro[15]}static/img/diviciones/{escudos_img[indice_escudo]}'
                img_left = Inches(left + (distancia_grupo / 2) - 0.25+ 0.85)
                slide.shapes.add_picture(escudo_path, img_left, top_img,
                                        width=img_width, height=img_height)
                indice_escudo += 1

            # Avanzar al siguiente grupo
            left += distancia_grupo + 1.65

    # Slide lamina de cuadros 
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    #slide.shapes.title.text = f"EVALUACIÓN LÍNEA DE OPERACIONES “PAYA” ".upper()
    titulo = f"""EVALUACIÓN LÍNEA DE OPERACIONES “PAYA” META MENSUAL "{mes_num_inicio}" """.upper() 
    agregar_texto(slide, titulo, 2, 0.2, font_size=28, color=[56, 87, 35], word_wrap= True, width_in=10, height_in=1, bold=True)
    agregar_texto(slide, fecha_actual, 8, 0.88)

    resultado = calcular_spoa.resultados_narcotrafico_boletin_metas()
        # Coordenadas de bloques de tabla
    bloques = [
        (36, 41,   0.7, 2.1),   # Fila 1 - izquierda
        (42, 47,   4.7, 2.1),  # Fila 1 - centro
        (48, 53,  8.75, 2.1),   # Fila 1 - derecha

        (54, 59,  0.7, 4.7),   # Fila 2 - izquierda
        (60, 65,  4.7, 4.7),  # Fila 2 - centro
        #(30, 35,  8.75, 4.7),   # Fila 2 - derecha

        #(36, 41,  0.9, 4.2),   # Fila 3 - izquierda
        #(42, 47,  5.05, 4.2),  # Fila 3 - centro
        #(48, 53,  9.2, 4.2),   # Fila 3 - derecha

        #(54, 59,  0.9, 5.7),   # Fila 4 - izquierda
        #(60, 65,  5.05, 5.7),  # Fila 4 - centro
        # (66, 71, 9.2, 5.7),  # Fila 4 - derecha (omitido)
    ]


    # Para cada bloque, pasa también las rutas de escudos
    for i, (inicio, fin, left, top) in enumerate(bloques):
        data = resultado[0][inicio:fin + 1]
        
        escudos_tabla = resultado[1][30 + i * 5 : 30 + (i + 1) * 5]

        crear_tabla_estilizada_con_escudos_5(prs, slide, data, 8, filtro[15], escudos_tabla, left=left, top=top, width=4)

    print(len(resultado[1]))
    
    # Ejemplo de uso:
    

    # === Lista de escudos ===
    escudos_img = [
        "div05_ayuda.png", "div06_ayuda.png", "div07_ayuda.png",
        "div08_ayuda.png", "omega_ayuda.png"
    ]

    # === Configuración general ===
    largo_linea = 1.5
    distancia_grupo = 2.35
    base_left = 0.7
    top_lineas = [1.7, 4.4]
    top_imagenes = [1.4, 4.1]
    img_width = Inches(0.5)
    img_height = Inches(0.5)

    # === Inserción automática ===
    indice_escudo = 0
    for fila in range(2):  # dos filas de líneas e imágenes
        top_linea = top_lineas[fila]
        top_img = Inches(top_imagenes[fila])
        left = base_left

        for _ in range(3 if fila == 0 else 2):  # Fila 1 = 3 bloques, Fila 2 = 2 bloques
            # Línea izquierda
            agregar_linea(slide, left_in=left, top_in=top_linea, largo_in=largo_linea,
                        orientacion="horizontal", color_rgb=(51, 63, 80), grosor_pt=3)

            # Línea derecha
            right_line = left + distancia_grupo
            agregar_linea(slide, left_in=right_line, top_in=top_linea, largo_in=largo_linea,
                        orientacion="horizontal", color_rgb=(51, 63, 80), grosor_pt=3)

            # Imagen centrada entre líneas
            if indice_escudo < len(escudos_img):
                escudo_path = f'{filtro[15]}static/img/diviciones/{escudos_img[indice_escudo]}'
                img_left = Inches(left + (distancia_grupo / 2) - 0.25+ 0.85)
                slide.shapes.add_picture(escudo_path, img_left, top_img,
                                        width=img_width, height=img_height)
                indice_escudo += 1

            # Avanzar al siguiente grupo
            left += distancia_grupo + 1.65

    
    datos, resultados_data, data_total = calcular_spoa.cocaina_en_proceso()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
 
    #slide.shapes.title.text = f"CONVERSIÓN POR CONCENTRACIÓN CLHC EN PROCESO ".upper()
    agregar_texto(slide, fecha_actual, 8, 0.88)
    color = [255,255,255, "Calibri", 14, True, False]
 
    cuadros_divisio_r(slide, anio_inicio , 1.85, 5.45, -90, color, 0.7)
    cuadros_divisio_r(slide, anio_inicio_f_f, 3.85, 5.45, -90, color, 0.7)

    color = [0,0,0, "Calibri", 14, True, False]
    cuadros_divisio_r(slide, datos[0] , 2.45, 5.3, 0, color, 1.2)
    cuadros_divisio_r(slide, datos[1] , 2.45, 5.6, 0, color, 1.2)

    cuadros_divisio_r(slide, datos[2] , 4.4, 5.3, 0, color, 1.2)
    cuadros_divisio_r(slide, datos[3] , 4.4, 5.6, 0, color, 1.2)

    cuadros_divisio_r(slide, datos[4] , 2.1, 6.22, 0, color, 1.65)
    cuadros_divisio_r(slide, datos[5] , 3.9, 6.22, 0, color, 1.65)

    cuadros_divisio_r(slide, datos[6] , 3.1, 6.87, 0, color, 1.4)

    color = [0,0,0, "Calibri", 14, True, False]
    
    # posiciones fijas de columnas
    x1 = 6.5
    x2 = 7.8

    # valores iniciales
    y_inicial = 2.3
    incremento = 0.45   # la distancia aproximada entre filas

    for i in range(0, len(resultados_data), 2):
        fila = i // 2
        y = y_inicial + (fila * incremento)

        if y == 6.8:
            color = [255,255,255, "Calibri", 14, True, False]

        cuadros_divisio_r(slide, resultados_data[i],   x1, y, 0, color, 1.15)
        cuadros_divisio_r(slide, resultados_data[i+1], x2, y, 0, color, 1.15)

    # posiciones fijas de columnas
    x1 = 10.75

    # valores iniciales
    y_inicial = 4.1
    incremento = 0.41   # distancia aproximada entre filas
    color_base = [0,0,0, "Calibri", 14, True, False]
    color_blanco = [255,255,255, "Calibri", 14, True, False]

    for i in range(len(data_total)):
        fila = i // 2
        y_inicial = y_inicial + incremento

        # Cambiar color en la fila 2
        if fila == 2:
            color = color_blanco
            y_inicial = y_inicial + incremento-0.25
            x1 = 9.4
            cuadros_divisio_r(slide, data_total[i], x1, y_inicial, 0, color, 2.5)
        else:
            color = color_base

            cuadros_divisio_r(slide, data_total[i], x1, y_inicial, 0, color, 1.15)

    y_inicial = y_inicial + incremento-0.1
    color = [128,0,0, "Calibri", 8, False, True]
    text = f"Meta {anio_inicio_f_f}"
    cuadros_divisio_r(slide, text, x1, y_inicial, 0, color, 2)

    color = [0,0,0, "Calibri", 7, False, True]
    text = f"Diferencia galones {anio_inicio} - {anio_inicio_f_f}"
    #cuadros_divisio_r(slide, text, 2.05, 6.45, 0, color, 1.7)
    text_1 = f"Diferencia kilogramos {anio_inicio} - {anio_inicio_f_f}"
    #cuadros_divisio_r(slide, text_1, 3.85, 6.45, 0, color, 1.7)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    #agregar_linea(slide, left_in=1, top_in=2, largo_in=3, orientacion="vertical", color_rgb=(0, 128, 0), grosor_pt=2)
    #agregar_linea(slide, left_in=1, top_in=3, largo_in=2, orientacion="diagonal", color_rgb=(0, 0, 255), grosor_pt=4)


    link_2 = link #carpeta base
    nombre_carpeta = f"DOCNA/"
    link = os.path.join(link_2, nombre_carpeta)
    dirercion_archvios = link
    
    if os.path.exists(link_2):
        for elemento in os.listdir(link_2):
            ruta = os.path.join(link_2, elemento)
            if os.path.isdir(ruta):
                shutil.rmtree(ruta)
            else:
                os.remove(ruta)
    else:
        os.mkdir(link_2)
    

    if not os.path.exists(link):
        os.mkdir(link)

    # Guardar archivo
    direcion = link+str("test")+'.pptx'
    prs.save(direcion)

    direcion = link+str("test")+'.pptx'

    DIRECION_3 = os.getenv('DIRECION_3')
    direcion= DIRECION_3+nombre_carpeta+str("test")+'.pptx'
    
    return [direcion, "nombre_doc"]