from flask import make_response
from datetime import date, time, datetime
from tipo_docker.k_c_boletin_comparativo_power_point.models.estadistica.estadistica_boletin_coe import *
from __init__ import *
from tipo_docker.z_z_configuarcion.caligrafia import *
from tipo_docker.z_z_configuarcion.fechas import *
from tipo_docker.z_z_configuarcion.header import *
from tipo_docker.z_z_configuarcion.logo import *
from tipo_docker.z_z_configuarcion.titulos import *  


 
from flask import make_response
from datetime import date, time, datetime



from tipo_docker.z_z_configuarcion.caligrafia import *
from tipo_docker.z_z_configuarcion.fechas import *
from tipo_docker.z_z_configuarcion.header import *
from tipo_docker.z_z_configuarcion.logo import *
from tipo_docker.z_z_configuarcion.titulos import *  
from tipo_docker.z_z_configuarcion.tipo_ayudas import *

from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.enum.chart import XL_TICK_MARK
from pptx.util import Pt

from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.chart import XL_TICK_LABEL_POSITION

from pptx.enum.shapes import MSO_SHAPE

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_VERTICAL_ANCHOR
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.dml import MSO_COLOR_TYPE
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_FILL
from pptx.oxml.xmlchemy import OxmlElement
import os
import shutil

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL



from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL


from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

import matplotlib.pyplot as plt
import numpy as np
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.dml.color import RGBColor


def grafico_reloj(ruta_salida, velocidad=0, color_reloj=(0,176,80), nombre='dirop'):
        """
        Genera un velocímetro tipo reloj y guarda la imagen PNG en la ruta indicada.
        Parámetros:
        - ruta_salida: ruta completa del archivo a guardar (sin extensión o con .png)
        - velocidad: valor entre 0 y 100
        """

        velocidad_max = 100
        velocidad = max(0, min(velocidad, velocidad_max))  # límite 0–100

        # --- Cálculo del ángulo ---
        angulo = np.deg2rad(-145 + (velocidad / velocidad_max) * 290)

        # --- Crear figura polar ---
        fig, ax = plt.subplots(figsize=(6, 6), subplot_kw={'projection': 'polar'})
        ax.set_theta_zero_location('N')
        ax.set_theta_direction(-1)
        ax.set_thetamin(-145)
        ax.set_thetamax(145)
        ax.set_ylim(0, 1)
        ax.set_yticklabels([])
        ax.set_xticklabels([])
        ax.grid(False)
        ax.spines['polar'].set_visible(False)

        # --- Fondo gris del arco ---
        theta_full = np.linspace(np.deg2rad(-145), np.deg2rad(145), 500)
        r_inner, r_outer = 0.1, 0.9
        for i in range(len(theta_full) - 1):
            ax.fill_between(
                [theta_full[i], theta_full[i + 1]],
                r_inner, r_outer,
                color="#D0D0D0",
                alpha=0.6
            )

        # --- Colorear según la velocidad ---
        def color_por_velocidad(v):
            if v < 35:
                return "#FF0000"     # rojo
            elif v < 50:
                return "#FFA600"     # naranja
            elif v < 80:
                return "#FFFF00"     # amarillo
            else:
                return "#00FF00"     # verde

        color = color_reloj

        theta_color = np.linspace(
            np.deg2rad(-145),
            np.deg2rad(-145 + (velocidad / velocidad_max) * 290),
            300
        )

        for i in range(len(theta_color) - 1):
            ax.fill_between(
                [theta_color[i], theta_color[i + 1]],
                r_inner, r_outer,
                color=color,
                alpha=0.9
            )

        # --- Líneas divisorias y numeración ---
        for v in range(0, velocidad_max + 1, 10):
            ang = np.deg2rad(-145 + (v / velocidad_max) * 290)
            ax.plot([ang, ang], [r_inner, r_outer], color='white', linewidth=2, zorder=6)

            # Texto de numeración
            r_text = r_outer + 0.08  # posición radial del texto
            ax.text(ang, r_text, str(v), ha='center', va='center',
                    color='gray', fontsize=25, fontweight='bold')

        # --- Aguja ---
        ax.plot([angulo, angulo], [0, 0.85], color='black', linewidth=3, zorder=10)
        ax.scatter(0, 0, color='black', s=120, zorder=11)


        # --- Fondo y guardado ---
        fig.patch.set_facecolor("#00000000")
        ax.set_facecolor("#404040")

        # Asegurar extensión PNG con nombre fijo
        nombre = nombre + '.png'
        if not ruta_salida.lower().endswith(nombre):
            if not ruta_salida.endswith(("/", "\\")):
                ruta_salida += "_"
            ruta_salida += nombre

        plt.savefig(ruta_salida, dpi=100, bbox_inches='tight', transparent=True)
        plt.close(fig)

        return ruta_salida

def set_cell_border(cell):
    """Aplica bordes negros a una celda de tabla en pptx."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for side in ['L', 'R', 'T', 'B']:
        ln = OxmlElement(f'a:ln{side}')
        ln.set('w', '12700')  # 1pt en EMUs

        solidFill = OxmlElement('a:solidFill')
        srgbClr = OxmlElement('a:srgbClr')
        srgbClr.set('val', '000000')
        solidFill.append(srgbClr)
        ln.append(solidFill)

        prstDash = OxmlElement('a:prstDash')
        prstDash.set('val', 'solid')
        ln.append(prstDash)

        tcPr.append(ln)

def apply_border(cell, color="000000", width=1):

    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    # 🔹 eliminar bordes existentes
    for child in list(tcPr):
        if child.tag.endswith(('lnL', 'lnR', 'lnT', 'lnB')):
            tcPr.remove(child)

    for line in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:

        ln = OxmlElement(line)
        ln.set('w', str(width * 12700))

        solidFill = OxmlElement('a:solidFill')
        srgbClr = OxmlElement('a:srgbClr')
        srgbClr.set('val', color)

        solidFill.append(srgbClr)
        ln.append(solidFill)

        tcPr.append(ln)

def compactar_texto_celda(cell, texto, font_name="Calibri", size_pt=9, bold=True, align=PP_ALIGN.CENTER):

    tf = cell.text_frame
    tf.clear()

    # quitar padding
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.margin_left = 0
    tf.margin_right = 0

    tf.word_wrap = False
    tf.auto_size = None
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.text = texto

    # 🔹 eliminar espacios extra
    p.space_before = Pt(0)
    p.space_after = Pt(0)

    # 🔹 interlineado simple
    p.line_spacing = 0.58

    # alineación horizontal
    p.alignment = align

    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold


def crear_tabla_afectacion(
    prs,
    slide,
    data,
    size_pt=9,
    left=0.5,
    top=1.2,
    width=None,
    font_name="Calibri",
    row_height_cm=0.5,
    header_height_cm=None,
    col_widths=None,
    header_color=RGBColor(127,127,127),
    indicadore=0,
    encabezado=True,
    ruta_base=""
):

    rows = len(data)
    cols = len(data[0])

    if width is None:
        width = sum(col_widths)

    row_height = Inches(row_height_cm / 2.54)
    header_height = Inches(header_height_cm / 2.54) if header_height_cm else row_height

    table_shape = slide.shapes.add_table(
        rows, cols,
        Inches(left), Inches(top),
        width,
        row_height * rows
    )

    table = table_shape.table

    # alturas
    for i in range(rows):
        table.rows[i].height = header_height if i == 0 else row_height

    # anchos
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # 🔥 PRECALCULO POSICIONES (CLAVE)
    col_offsets = [0]
    for c in range(1, cols):
        col_offsets.append(col_offsets[c-1] + table.columns[c-1].width)

    row_offsets = [0]
    for r in range(1, rows):
        row_offsets.append(row_offsets[r-1] + table.rows[r-1].height)

    table_left = table_shape.left
    table_top = table_shape.top

    fila_inicio = 1 if encabezado else 0

    # colores
    white = RGBColor(255,255,255)
    black = RGBColor(0,0,0)
    fila_clara = RGBColor(236,236,226)

    rojo = RGBColor(192,0,0)
    verde = RGBColor(0,176,80)

    for i, row in enumerate(data):
        for j, val in enumerate(row):

            cell = table.cell(i, j)

            # ----------------------------------------
            # 🔹 ALINEACIÓN (RESPETADA)
            # ----------------------------------------
            if indicadore == 0:
                if i == 0:
                    alineacion = PP_ALIGN.CENTER
                elif j == 2:
                    alineacion = PP_ALIGN.LEFT
                else:
                    alineacion = PP_ALIGN.CENTER
            else:
                if i == 0:
                    alineacion = PP_ALIGN.CENTER
                elif j == 0:
                    alineacion = PP_ALIGN.LEFT
                elif j == 4:
                    alineacion = PP_ALIGN.RIGHT
                else:
                    alineacion = PP_ALIGN.CENTER

            if encabezado == False:
                if j == 0:
                    alineacion = PP_ALIGN.LEFT
                elif j == 4:
                    alineacion = PP_ALIGN.RIGHT
                else:
                    alineacion = PP_ALIGN.CENTER

            # ----------------------------------------
            # 🔹 TEXTO + REDONDEO
            # ----------------------------------------
            texto = str(val)

            if i >= fila_inicio and j != 0:
                try:
                    texto = str(formato_numero(val))
                except:
                    pass

            compactar_texto_celda(
                cell,
                texto,
                font_name=font_name,
                size_pt=size_pt,
                bold=True,
                align=alineacion
            )

            run = cell.text_frame.paragraphs[0].runs[0]

            # ----------------------------------------
            # 🔹 COLORES BASE
            # ----------------------------------------
            cell.fill.solid()
            cell.fill.fore_color.rgb = fila_clara
            run.font.color.rgb = black

            # encabezado
            if encabezado and i == 0:
                cell.fill.fore_color.rgb = header_color
                run.font.color.rgb = white

            # ----------------------------------------
            # BORDES
            # ----------------------------------------
            apply_border(cell)

    return table

def cuadros_divisio(slide, valor_division, x, y):
  # ... dentro de tu función o flujo
    shapes = slide.shapes
    left = Inches(x)
    top = Inches(y)

    width = Inches(1.1)
    height = Inches(0.3)

    shape = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )

    # Eliminar fondo del rectángulo
    shape.fill.background()

    # Eliminar bordes
    shape.line.fill.background()

    # Agregar texto
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = str(valor_division)

    # Estilo de texto
    font = run.font
    font.name = 'Arial'
    font.size = Pt(14)
    font.bold = True
    font.color.rgb = RGBColor(0, 0, 0)

    # Alineación opcional
    p.alignment = PP_ALIGN.LEFT

def cuadros_divisio_r(slide, valor_division, x, y, rotar, color, ancho):
    shapes = slide.shapes
    left = Inches(x)
    top = Inches(y)
    width = Inches(ancho)
    height = Inches(0.3)

    shape = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )

    # Eliminar fondo y bordes
    shape.fill.background()
    shape.line.fill.background()

    # Agregar texto
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = str(valor_division)

    # Estilo de texto
    font = run.font
    font.name = color[3]
    font.size = Pt(color[4])
    font.bold = color[5]
    font.italic = color[6]
    font.color.rgb = RGBColor(color[0],color[1],color[2])

    # Alineación
    p.alignment = PP_ALIGN.CENTER

    # 🔄 Rotar el texto o el shape (90 grados)
    shape.rotation = rotar  # rota todo el cuadro de texto

    # Si quieres solo el texto en vertical sin girar la forma:
    # text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    # text_frame.text_anchor = MSO_VERTICAL_ANCHOR.VERTICAL

# Función para agregar línea personalizada
def agregar_linea(slide, left_in, top_in, largo_in, orientacion="horizontal", color_rgb=(0, 0, 0), grosor_pt=2):
    left = Inches(left_in)
    top = Inches(top_in)

    if orientacion == "horizontal":
        width = Inches(largo_in)
        height = Inches(0)  # Línea horizontal
    elif orientacion == "vertical":
        width = Inches(0)
        height = Inches(largo_in)  # Línea vertical
    elif orientacion == "diagonal":
        width = Inches(largo_in)
        height = Inches(largo_in)  # Diagonal con 45°
    else:
        raise ValueError("Orientación no válida. Usa: horizontal, vertical o diagonal.")

    line = slide.shapes.add_shape(
        autoshape_type_id=1,  # Línea
        left=left,
        top=top,
        width=width,
        height=height
    )

    # Estilizar línea
    line.line.color.rgb = RGBColor(*color_rgb)
    line.line.width = Pt(grosor_pt)

def merge_repetidos_columna(tabla, data, col_index=1,
                            font_name="Calibri",
                            font_size=11,
                            color=(0,0,0),
                            bold=True):

    # -----------------------------
    # VALIDAR DATA
    # -----------------------------
    if not data:
        return

    # eliminar filas vacías
    data = [row for row in data if row]

    # si después de limpiar queda menos de 2 filas
    if len(data) < 2:
        return

    # verificar columnas
    for row in data:
        if len(row) <= col_index:
            return

    start = 1
    actual = data[start][col_index]

    for i in range(start + 1, len(data)):

        if data[i][col_index] != actual:

            end = i - 1

            if end > start:
                merge_clean(
                    tabla,
                    start, col_index,
                    end, col_index,
                    text=actual,
                    font_name=font_name,
                    font_size=font_size,
                    color=color,
                    bold=bold
                )

            start = i
            actual = data[i][col_index]

    # -----------------------------
    # último grupo
    # -----------------------------
    end = len(data) - 1
    if end > start:
        merge_clean(
            tabla,
            start, col_index,
            end, col_index,
            text=actual,
            font_name=font_name,
            font_size=font_size,
            color=color,
            bold=bold
        )

def merge_clean(
        table,
        r1, c1,
        r2, c2,
        text=None,
        font_name="Calibri",
        font_size=9,
        color=(0,0,0),
        bold=True,
        align="center"
    ):

    # ----------------------------------------
    # limpiar celdas
    # ----------------------------------------
    for r in range(r1, r2+1):
        for c in range(c1, c2+1):
            if not (r == r1 and c == c1):
                table.cell(r,c).text = ""

    # ----------------------------------------
    # fusionar
    # ----------------------------------------
    try:
        table.cell(r1,c1).merge(table.cell(r2,c2))
    except:
        pass

    cell = table.cell(r1,c1)

    # ----------------------------------------
    # centrado vertical
    # ----------------------------------------
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcPr.set("anchor", "ctr")

    # ----------------------------------------
    # limpiar text frame
    # ----------------------------------------
    tf = cell.text_frame
    tf.clear()

    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.margin_left = 0
    tf.margin_right = 0

    # ----------------------------------------
    # texto
    # ----------------------------------------
    p = tf.paragraphs[0]
    p.text = text if text else ""

    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1

    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT

    # ----------------------------------------
    # fuente
    # ----------------------------------------
    if not p.runs:
        run = p.add_run()
    else:
        run = p.runs[0]

    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)

    return cell

def agregar_texto(
    slide,
    texto,
    left_in,
    top_in,
    width_in=0.5,
    height_in=0.3,
    font_size=12,
    font_name="Calibri",
    color=(127,127,127),
    bold=False,
    italic=False,
    align="left",
    word_wrap=False,

    # 🔹 NUEVOS PARÁMETROS
    fill=False,
    fill_color=(255,255,255),

    margin_left=0,
    margin_right=0,
    margin_top=0,
    margin_bottom=0,
    underline=False, # 👈 subrayado
    strike=False, # 👈 tachado
    shadow=False
):
    """
    Inserta un textbox en PowerPoint.

    Por defecto:
    - Sin relleno
    - Sin márgenes
    - Texto alineado según parámetro

    Parámetros nuevos:
    fill -> activar relleno
    fill_color -> color del relleno
    margin_* -> márgenes internos
    """

    # -------------------------------------------------
    # POSICIÓN Y TAMAÑO
    # -------------------------------------------------
    left = Inches(left_in)
    top = Inches(top_in)
    width = Inches(width_in)
    height = Inches(height_in)

    textbox = slide.shapes.add_textbox(left, top, width, height)

    # -------------------------------------------------
    # CONFIGURAR RELLENO
    # -------------------------------------------------
    if fill:
        textbox.fill.solid()
        textbox.fill.fore_color.rgb = RGBColor(*fill_color)
    else:
        textbox.fill.background()   # sin relleno

    # -------------------------------------------------
    # TEXTO
    # -------------------------------------------------
    tf = textbox.text_frame

    tf.word_wrap = word_wrap
    tf.auto_size = None
    tf.clear()

    # -------------------------------------------------
    # MÁRGENES
    # -------------------------------------------------
    tf.margin_left = Pt(margin_left)
    tf.margin_right = Pt(margin_right)
    tf.margin_top = Pt(margin_top)
    tf.margin_bottom = Pt(margin_bottom)

    # -------------------------------------------------
    # PÁRRAFO
    # -------------------------------------------------
    p = tf.paragraphs[0]
    p.text = texto

    # alineación
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    elif align == "justify":
        p.alignment = PP_ALIGN.JUSTIFY
    else:
        p.alignment = PP_ALIGN.LEFT

    # -------------------------------------------------
    # FUENTE
    # -------------------------------------------------
    font = p.font
    font.size = Pt(font_size)
    font.name = font_name
    font.bold = bold
    font.italic = italic
    font.color.rgb = RGBColor(*color)
    font.underline = underline   # 👈 subrayado
    font.strike = strike         # 👈 tachado
    font.shadow = shadow

    return textbox


def agregar_cuadro_redondo(
    slide,
    texto,
    left_in,
    top_in,
    width_in=0.5,
    height_in=0.3,
    font_size=12,
    font_name="Calibri",
    color=(0,0,0),

    bold=False,
    italic=False,
    align="left",

    fill=True,
    fill_color=(255,255,255),

    border=True,
    border_color=(0,0,0),
    border_width=1,

    margin_left=5,
    margin_right=5,
    margin_top=3,
    margin_bottom=3
):
    
    # -------------------------------------------------
    # POSICIÓN
    # -------------------------------------------------
    left = Inches(left_in)
    top = Inches(top_in)
    width = Inches(width_in)
    height = Inches(height_in)

    # -------------------------------------------------
    # CREAR RECTÁNGULO REDONDO
    # -------------------------------------------------
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height
    )

    # -------------------------------------------------
    # RELLENO
    # -------------------------------------------------
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill_color)
    else:
        shape.fill.background()

    # -------------------------------------------------
    # BORDE
    # -------------------------------------------------
    if border:
        shape.line.color.rgb = RGBColor(*border_color)
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()

    # -------------------------------------------------
    # TEXTO
    # -------------------------------------------------
    tf = shape.text_frame
    tf.clear()

    tf.margin_left = Pt(margin_left)
    tf.margin_right = Pt(margin_right)
    tf.margin_top = Pt(margin_top)
    tf.margin_bottom = Pt(margin_bottom)

    p = tf.paragraphs[0]
    p.text = texto

    # alineación
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    elif align == "justify":
        p.alignment = PP_ALIGN.JUSTIFY
    else:
        p.alignment = PP_ALIGN.LEFT

    # -------------------------------------------------
    # FUENTE
    # -------------------------------------------------
    font = p.font
    font.size = Pt(font_size)
    font.name = font_name
    font.bold = bold
    font.italic = italic
    font.color.rgb = RGBColor(*color)

    return shape

def colorear_celdas_unidas(tabla, data, col_index=1):

    # validar que exista contenido
    if not data or len(data) < 2:
        return

    # validar columnas
    for row in data:
        if len(row) <= col_index:
            return

    color1 = RGBColor(242,242,242)
    color2 = RGBColor(217,217,217)

    color_actual = color1
    start = 1
    actual = data[start][col_index]

    for i in range(2, len(data)):

        if data[i][col_index] != actual:

            cell = tabla.cell(start, col_index)
            cell.fill.solid()
            cell.fill.fore_color.rgb = color_actual

            # alternar color
            color_actual = color2 if color_actual == color1 else color1

            start = i
            actual = data[i][col_index]

    # último grupo
    cell = tabla.cell(start, col_index)
    cell.fill.solid()
    cell.fill.fore_color.rgb = color_actual

def reloj_grafico(slide_cuerpo_5, reloj_base, altura, izquierda, img_path, texto, colo_texto, color_relleno):
    slide_cuerpo_5.shapes.add_picture(
    reloj_base,
    Inches(altura),   # izquierda
    Inches(izquierda),     # arriba
    width=Inches(3),
    height=Inches(3)
    )
    altura_1 = altura + 0.6
    izquierda_1 = izquierda + 0.52
    slide_cuerpo_5.shapes.add_picture(
    img_path,
    Inches(altura_1),   # izquierda
    Inches(izquierda_1),     # arriba
    width=Inches(1.8),
    height=Inches(1.8)
    )
    altura = altura + 0.9 
    izquierda = izquierda + 2.6 

    agregar_cuadro_redondo(
        slide_cuerpo_5,
        texto=texto,
        left_in=altura,
        top_in=izquierda,
        width_in=1.3,
        height_in=0.4,
        font_size=14,
        bold=True,
        align="center",
        fill=True,
        color=colo_texto,
        fill_color=color_relleno,
        border_color=color_relleno,
        border_width=2
    )

def tringulos(lamina_comparitva, nun_x, datos, ruta_base, left, suma, encabezado=False):

    for i, x in enumerate(datos):
        if encabezado == True:
            if i == 0:
                continue  # saltar encabezado

        valor = float(x[4])

        if valor > 0:
            trinagulo_color = "tri_verde_new.png"
        elif valor < 0:
            trinagulo_color = "tri_rojo_new.png"
        else:
            trinagulo_color = "tri_amarillo_new.png"

        tirangulo = f"{ruta_base}/static/img/{trinagulo_color}"

        lamina_comparitva.shapes.add_picture(
            tirangulo,
            Inches(left),
            Inches(nun_x),
            width=Inches(0.14),
            height=Inches(0.14)  # 🔥 cuadrado fijo
        )

        nun_x += suma
#comparativo por divisiones

def comparativo_comparativo_mapa(
    fecha_inicial_u_l, fecha_final_u_l,
    fecha_inicial_p_l, fecha_final_p_l,
    filtro, link, puerto
    ):

    # 🔹 Fechas periodo actual
    f_ini = fecha(fecha_inicial_u_l)
    f_fin = fecha(fecha_final_u_l)

    # 🔹 Fechas periodo comparación
    f_ini_comp = fecha(fecha_inicial_p_l)
    f_fin_comp = fecha(fecha_final_p_l)

    # 🔹 Desglose útil
    dia_ini, mes_ini, anio_ini = f_ini[0], f_ini[1], f_ini[2]
    dia_fin, mes_fin, anio_fin = f_fin[0], f_fin[1], f_fin[2]

    dia_ini_c, mes_ini_c, anio_ini_c, mes_ini_num = f_ini_comp[0], f_ini_comp[1], f_ini_comp[2], f_ini_comp[3]
    dia_fin_c, mes_fin_c, anio_fin_c, mes_fin_num = f_fin_comp[0], f_fin_comp[1], f_fin_comp[2], f_fin_comp[3]





    # 👉 Aquí continúa tu lógica (mapas, consultas, etc.)

     # 👉 Lógica (mapas, consultas, etc.)
    resultados_spoa = Calculo_Spoa()

    datos_infreso = resultados_spoa.comparativo_mapa(
        fecha_inicial_u_l,
        fecha_final_u_l,
        fecha_inicial_p_l,
        fecha_final_p_l,
        filtro,
    )

    # 🔹 Texto descriptivo (ANTES vs AHORA)
    fechas = (
        f"{dia_ini} DE {mes_ini} hasta {dia_fin} de {mes_fin} del {anio_fin_c} - {anio_ini}"
    )

    fechas_2 = (
        f"{dia_ini} DE {mes_ini} hasta {dia_fin} de {mes_fin} del {anio_ini}"
    )
    
    fechas = fechas.upper()


    titulo=f"RESULTADOS OPERACIONALES COMPARATIVO  \n{anio_ini_c} - {anio_fin_c}"

    ruta_base = filtro[15]
    ayuda = f"{ruta_base}/static/template/template_ayuda_dirop_2026.pptx"
    prs = Presentation(ayuda)
    #este codigo es para saber que tipo de ayudas tengo
    #tipo_ayudas_template(prs)
   
    slide_portada = prs.slides.add_slide(prs.slide_layouts[4])
    agregar_texto(slide_portada, titulo, 7.715, 2.515, font_size=32, color=[180,180,180], word_wrap= True, width_in=4, height_in=2, bold=True, align =  '', font_name="Calibri", shadow=True)
    agregar_texto(slide_portada, titulo, 7.7, 2.5, font_size=32, color=[154,0,0], word_wrap= True, width_in=4, height_in=2, bold=True, align =  '', font_name="Calibri", shadow=True)

    agregar_texto(slide_portada, "EJÉRCITO NACIONAL", 7.7, 4.6, font_size=20, color=[89,89,89], word_wrap= True, width_in=4, height_in=0.5, bold=True, align =  '', font_name="Calibri", shadow=True)

    lamina_comparitva = prs.slides.add_slide(prs.slide_layouts[2])



    
    titulo = f"""RESULTADOS OPERACIONALES \nCOMPARATIVO {fechas}"""

    agregar_texto(lamina_comparitva, titulo, 1.7, 0.2, font_size=25, color=[192,0,0], word_wrap= True, width_in=10, height_in=1, bold=True, align =  'center', font_name="Calibri", shadow=True)
    
    titulo = f"""{fechas_2}""".upper()

    agregar_texto(lamina_comparitva, titulo, 8.8, 1.02, font_size=11, color=[140,140,140], word_wrap= True, width_in=4, height_in=0.2, bold=True, align =  'right', font_name="Calibri", shadow=True)

    agregar_texto(lamina_comparitva, "AFECTACIÓN A LA CAPACIDAD ARMADA DE LA AMENAZA", 0.5, 1.2, font_size=11, color=[255,255,255], word_wrap= True, width_in=6, height_in=0.2, bold=True, align = 'center', fill=True, fill_color=(94,119,89), font_name="Calibri"  )

    col_widths = [

        Inches(2.8),
        Inches(0.8),
        Inches(0.8),
        Inches(0.8),
        Inches(0.8),
        ]

    header_color = RGBColor(94,119,89)
    row_height_cm=0.6
    datos =datos_infreso[0]
    crear_tabla_afectacion(
            prs,
            lamina_comparitva,
            datos,
            size_pt=11,
            left=0.5,
            top=1.4,
            width=9,
            font_name="Calibri",
            row_height_cm=0.5,
            header_height_cm=row_height_cm,
            col_widths = col_widths,
            header_color = header_color,
            indicadore = 1,
            ruta_base=ruta_base
    )
    tringulos(lamina_comparitva, 1.66, datos, ruta_base, 5.85, 0.21, encabezado=True )



    agregar_texto(lamina_comparitva, "AFECTACIÓN A LA CAPACIDAD ARMADA DE LA AMENAZA", 0.5, 2.75, font_size=11, color=[255,255,255], word_wrap= True, width_in=6, height_in=0.2, bold=True, align = 'center', fill=True, fill_color=(94,119,89), font_name="Calibri"  )

    datos =datos_infreso[1]
  

    crear_tabla_afectacion(
            prs,
            lamina_comparitva,
            datos,
            size_pt=11,
            left=0.5,
            top=2.95,
            width=9,
            font_name="Calibri",
            row_height_cm=0.5,
            header_height_cm=row_height_cm,
            col_widths = col_widths,
            header_color = header_color,
            indicadore = 1,
            encabezado = False,
            ruta_base=ruta_base
    )
    tringulos(lamina_comparitva, 2.96, datos, ruta_base, 5.85, 0.21, encabezado=False )

    agregar_texto(lamina_comparitva, "COMBATES", 0.5, 3.85, font_size=11, color=[255,255,255], word_wrap= True, width_in=6, height_in=0.2, bold=True, align = 'center', fill=True, fill_color=(94,119,89), font_name="Calibri"  )

    datos =datos_infreso[2]

    crear_tabla_afectacion(
            prs,
            lamina_comparitva,
            datos,
            size_pt=11,
            left=0.5,
            top=4.05,
            width=9,
            font_name="Calibri",
            row_height_cm=0.5,
            header_height_cm=row_height_cm,
            col_widths = col_widths,
            header_color = header_color,
            indicadore = 1,
            encabezado = False,
            ruta_base=ruta_base
    )
    tringulos(lamina_comparitva, 4.06, datos, ruta_base, 5.85, 0.21, encabezado=False )

    agregar_texto(lamina_comparitva, "GUERRA DE MINAS", 0.5, 4.7, font_size=11, color=[255,255,255], word_wrap= True, width_in=6, height_in=0.2, bold=True, align = 'center', fill=True, fill_color=(94,119,89), font_name="Calibri"  )

    datos =datos_infreso[3]

    crear_tabla_afectacion(
            prs,
            lamina_comparitva,
            datos,
            size_pt=11,
            left=0.5,
            top=4.9,
            width=9,
            font_name="Calibri",
            row_height_cm=0.5,
            header_height_cm=row_height_cm,
            col_widths = col_widths,
            header_color = header_color,
            indicadore = 1,
            encabezado = False,
            ruta_base=ruta_base
    )
    tringulos(lamina_comparitva, 4.95, datos, ruta_base, 5.85, 0.21, encabezado=False )
    altura=0.4
    agregar_texto(lamina_comparitva, "AFECTACIÓN A LAS ECONOMÍAS ILÍCITAS – NARCOTRÁFICO", 6.8, 1.2, font_size=11, color=[255,255,255], word_wrap= True, width_in=6, height_in=0.2, bold=True, align = 'center', fill=True, fill_color=(94,119,89), font_name="Calibri"  )

    datos =datos_infreso[4]

    crear_tabla_afectacion(prs, lamina_comparitva, datos, size_pt=11, left=6.8, top=1.4, width=9, row_height_cm=0.4, col_widths = col_widths, header_color = header_color, indicadore = 1, ruta_base=ruta_base)
    tringulos(lamina_comparitva, 1.65, datos, ruta_base, 12.15, 0.20, encabezado=True )

    agregar_texto(lamina_comparitva, "AFECTACIÓN A LAS ECONOMÍAS ILÍCITAS – EIYM", 6.8, 2.9, font_size=11, color=[255,255,255], word_wrap= True, width_in=6, height_in=0.2, bold=True, align = 'center', fill=True, fill_color=(94,119,89), font_name="Calibri"  )

    datos =datos_infreso[5]

    crear_tabla_afectacion(
            prs,
            lamina_comparitva,
            datos,
            size_pt=11,
            left=6.81,
            top=3.1,
            width=9,
            font_name="Calibri",
            header_height_cm=row_height_cm,
            col_widths = col_widths,
            header_color = header_color,
            indicadore = 1,
             encabezado = False,
             ruta_base=ruta_base
    )
    tringulos(lamina_comparitva, 3.15, datos, ruta_base, 12.15, 0.21, encabezado=False )
    
    agregar_texto(lamina_comparitva, "AFECTACIÓN A LAS ECONOMÍAS ILÍCITAS - CONTRA LA LIBERTAD PERSONAL", 6.8, 4, font_size=11, color=[255,255,255], word_wrap= True, width_in=6, height_in=0.2, bold=True, align = 'center', fill=True, fill_color=(94,119,89), font_name="Calibri"  )

    datos =datos_infreso[6]

    crear_tabla_afectacion(
            prs,
            lamina_comparitva,
            datos,
            size_pt=11,
            left=6.8,
            top=4.2,
            width=9,
            font_name="Calibri",
            row_height_cm=altura,
            header_height_cm=row_height_cm,
            col_widths = col_widths,
            header_color = header_color,
            indicadore = 1,
             encabezado = False,
             ruta_base=ruta_base
    )

    tringulos(lamina_comparitva, 4.22, datos, ruta_base, 12.15, 0.25, encabezado=False )
    
    agregar_texto(lamina_comparitva, "AFECTACIÓN A LAS ECONOMÍAS ILÍCITAS - HIDROCARBUROS", 6.8, 4.65, font_size=11, color=[255,255,255], word_wrap= True, width_in=6, height_in=0.2, bold=True, align = 'center', fill=True, fill_color=(94,119,89), font_name="Calibri"  )

    datos =datos_infreso[7]

    crear_tabla_afectacion(
            prs,
            lamina_comparitva,
            datos,
            size_pt=11,
            left=6.8,
            top=4.85,
            width=9,
            font_name="Calibri",
            row_height_cm=altura,
            header_height_cm=row_height_cm,
            col_widths = col_widths,
            header_color = header_color,
            indicadore = 1,
            encabezado = False,
            ruta_base=ruta_base
    )

    tringulos(lamina_comparitva, 4.87, datos, ruta_base, 12.15, 0.25, encabezado=False )
    agregar_texto(lamina_comparitva, "AFECTACIÓN A LAS ECONOMÍAS ILÍCITAS - CONTRABANDO", 6.8, 5.3, font_size=11, color=[255,255,255], word_wrap= True, width_in=6, height_in=0.2, bold=True, align = 'center', fill=True, fill_color=(94,119,89), font_name="Calibri"  )

    datos =datos_infreso[8]

    crear_tabla_afectacion(
            prs,
            lamina_comparitva,
            datos,
            size_pt=11,
            left=6.8,
            top=5.5,
            width=9,
            font_name="Calibri",
            row_height_cm=altura,
            header_height_cm=row_height_cm,
            col_widths = col_widths,
            header_color = header_color,
            indicadore = 1,
            encabezado = False,
            ruta_base=ruta_base
    )
    tringulos(lamina_comparitva, 5.52, datos, ruta_base, 12.15, 0.22, encabezado=False )
    
    agregar_texto(lamina_comparitva, "LOE AMAZONÍA", 4.1, 6.4, font_size=11, color=[255,255,255], word_wrap= True, width_in=6, height_in=0.2, bold=True, align = 'center', fill=True, fill_color=(94,119,89), font_name="Calibri"  )

    datos =datos_infreso[9]
    crear_tabla_afectacion(
            prs,
            lamina_comparitva,
            datos,
            size_pt=11,
            left=4.1,
            top=6.6,
            width=9,
            font_name="Calibri",
            row_height_cm=altura,
            header_height_cm=row_height_cm,
            col_widths = col_widths,
            header_color = header_color,
            indicadore = 1,
            encabezado = False,
            ruta_base=ruta_base
    )
    tringulos(lamina_comparitva, 6.62, datos, ruta_base, 9.45, 0.22, encabezado=False )
    link_2 = link #carpeta base
    nombre_carpeta = f"DIROP_COMP"
    link = os.path.join(link_2, nombre_carpeta)
    dirercion_archvios = link
    
    if os.path.exists(link_2):
        for elemento in os.listdir(link_2):
            ruta = os.path.join(link_2, elemento)
            if os.path.isdir(ruta):
                shutil.rmtree(ruta)
            else:
                os.remove(ruta)
    else:
        os.mkdir(link_2)
    

    if not os.path.exists(link):
        os.mkdir(link)

    # Guardar archivo
    direcion = link+str("test")+'.pptx'
    prs.save(direcion)

    direcion = link+str("test")+'.pptx'

    DIRECION_3 = os.getenv('DIRECION_3')
    direcion= DIRECION_3+nombre_carpeta+str("test")+'.pptx'
    


    nombre = f"RESULTADOS OPERACIONALES COMPARATIVO {dia_fin}-{mes_fin}-{anio_ini}".upper()
    return [direcion, nombre]