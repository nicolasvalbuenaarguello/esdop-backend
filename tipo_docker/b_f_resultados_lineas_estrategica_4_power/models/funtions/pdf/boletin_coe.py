from flask import make_response
from datetime import date, time, datetime
from tipo_docker.b_f_resultados_lineas_estrategica_4_power.models.estadistica.estadistica_boletin_coe import *
from __init__ import *
from tipo_docker.z_z_configuarcion.caligrafia import *
from tipo_docker.z_z_configuarcion.fechas import *
from tipo_docker.z_z_configuarcion.header import *
from tipo_docker.z_z_configuarcion.logo import *
from tipo_docker.z_z_configuarcion.titulos import *  
from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.enum.chart import XL_TICK_MARK
from pptx.util import Pt

from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.chart import XL_TICK_LABEL_POSITION

from pptx.enum.shapes import MSO_SHAPE

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.dml import MSO_COLOR_TYPE
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_FILL
from pptx.oxml.xmlchemy import OxmlElement


def cuadros_divisio(slide, valor_division, x, y):
    shapes = slide.shapes
    left = Inches(x)
    top = Inches(y)

    width = Inches(0.6)
    height = Inches(0.2)

    shape = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )

    # subtitle = slide.placeholders[1]
    text_frame = shape.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = str(valor_division)

    font = run.font
    font.name = 'Calibri'
    font.size = Pt(11)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_1


    # shape.text = text_frame
    # font = shape.font
    # p.alignment = PP_ALIGN.LEFT
    font.color.rgb = RGBColor(255, 255, 255)

    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(20, 62, 52)


def charts(slide, num_placeholder, categories, series_ad, series_nombre, afe):


    placeholder = slide.placeholders[num_placeholder]
    chart_data = ChartData()
    chart_data.categories = categories
    chart_data.add_series(series_nombre,series_ad)
    
    
    graphic_frame = placeholder.insert_chart(XL_CHART_TYPE.DOUGHNUT, chart_data)
    chart = graphic_frame.chart
    chart.chart_type

    chart.has_legend = True
    chart.font.size = Pt(14)
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.font.color.rgb = RGBColor(70, 70, 70)



    categorias = chart.plots[0]
    categorias.has_text_frame = False

    plot = chart.plots[0]
    plot.has_data_labels = True

    data_labels = plot.data_labels

    data_labels.font.size = Pt(14)
    data_labels.font.color.rgb = RGBColor(125, 0, 0)
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    
    chart.has_title = False

    colors = {
        series_nombre: ['3fb872', '035827', '0d2116'],  # Example colors
        # Add more series names and color lists as needed
    }
    for series in chart.series:
        if series.name in colors:
            for i, color_code in enumerate(colors[series.name]):
                if i < len(series.points):
                    fill = series.points[i].format.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor.from_string(color_code)
        i = 0
        series = chart.series[0]
        for point in series.points:
            point.data_label.has_text_frame = True
            #Assigning custom text for data label associated with each data-point
            point.data_label.text_frame.text =  str(int(chart.series[0].values[i])) +" "+afe
       
            for run in point.data_label.text_frame.paragraphs[0].runs:
                run.font.size = Pt(12)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0, 0, 0)
            i+=1
    
def tabla(slide, datos, anio_ant, anio_act, num_placeholder, afectacion):
    
    tamnio = len(datos)

    placeholder = slide.placeholders[num_placeholder]
    graphic_frame = placeholder.insert_table(rows=tamnio+2, cols=6)
    table = graphic_frame.table

    start_cell = table.cell(0, 0)
    end_cell = table.cell(0,5)
    start_cell.merge(end_cell)

    cell = table.cell(0, 0)
    cell.text = afectacion
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.size = Pt(12)
    paragraph.font.bold = True

    cell = table.cell(1, 0)
    cell.text = 'DIVISIÓN'
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.size = Pt(12)
    paragraph.font.bold = True
  
    cell = table.cell(1, 1)
    cell.text = str(anio_ant)
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.size = Pt(12)
    paragraph.font.bold = True

    cell = table.cell(1, 2)
    cell.text = str(anio_act)
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.size = Pt(12)
    paragraph.font.bold = True

    cell = table.cell(1, 3)
    cell.text = str("DIF.")
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.size = Pt(12)
    paragraph.font.bold = True
        
    cell = table.cell(1, 4)
    cell.text = str("%")
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.size = Pt(12)
    paragraph.font.bold = True
            
    cell = table.cell(1, 5)
    cell.text = str("SEM.")
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.size = Pt(12)
    paragraph.font.bold = True

    table.columns[0].width = Inches(1)
    table.columns[1].width = Inches(0.5)
    table.columns[2].width = Inches(0.6)
    table.columns[3].width = Inches(0.6)
    table.columns[4].width = Inches(0.6)
    table.columns[5].width = Inches(0.6)


    numero = 1
    fila = 2
    
    tamnio = len(datos)
    for x in datos:   
            if fila <= tamnio+1:
                row = table.rows[fila]
                row.height = Inches(0.1)

                cols = table.columns[1]
                cols.width = Inches(1)

                cell = table.cell(fila, 0)
                cell.text = x[0]
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(9)
                paragraph.font.bold = True

                cell = table.cell(fila, 1)
                cell.text = str(x[1])
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(9)
                paragraph.font.bold = True
                
                cell = table.cell(fila, 2)
                cell.text = str(x[2])
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(9)
                paragraph.font.bold = True
                                
                cell = table.cell(fila, 3)
                cell.text = str(x[3])
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(9)
                paragraph.font.bold = True
                
                cell = table.cell(fila, 4)
                cell.text = str(x[4])
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(9)
                paragraph.font.bold = True
                                
                cell = table.cell(fila, 5)
                cell.text = str(x[5])
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(9)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(255, 255, 255) 

                fila= fila+1

    r=2
    for x in datos:
        c=5
        for y in x:
            
            #Insert data into table
            if x[5] == "-":
                cell = table.cell(r, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(128, 0, 0)    
            elif x[5] == "+":
                cell = table.cell(r, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(55, 86, 35)   
            else:
                cell = table.cell(r, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(255, 192, 0)     
            cell = apply_border(cell, edges = ["left", "right", "top", "bottom"])

        r = r +1


def apply_border(cell, edges = ["left", "right", "top", "bottom"], border_color="000000", border_width=1):
    if type(edges) is not list: edges = [edges]
    border_width = str(border_width*Pt(1))
    def SubElement(parent, tagname, **kwargs):
            element = OxmlElement(tagname)
            element.attrib.update(kwargs)
            parent.append(element)
            return element
    
    lines = [{"left": 'a:lnL',
              "right": 'a:lnR',
              "top": 'a:lnT',
              "bottom": 'a:lnB'}[_] for _ in edges]
    
    if cell.fill.type == MSO_FILL.SOLID: fill_color = cell.fill.fore_color.rgb
    cell.fill.background()
    
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for line in lines:
        
        # Remove duplicate tag if it exists
        tag = line.split(":")[-1]
        for e in tcPr.getchildren():
            if tag in str(e.tag): tcPr.remove(e)
        
        ln = SubElement(tcPr, line, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(ln, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(ln, 'a:prstDash', val='solid')
        round_ = SubElement(ln, 'a:round')
        headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')
        
    cell.fill.solid()
    if cell.fill.type == MSO_FILL.SOLID: cell.fill.fore_color.rgb = fill_color
    return(cell)

#comparativo por divisiones
def comparativo_comparativo_mapa(fecha_inicial_u_l, fecha_final_u_l, fecha_inicial_p_l, fecha_final_p_l, filtro, dirercion_archvios, nombre_carpeta):

    #-----------------------------------
    #modulo para el calculo de los resultados 
    fechas_inicial = fecha(fecha_inicial_u_l)
    fechas_final = fecha(fecha_final_u_l)

    fechas_inicial_final = fecha(fecha_inicial_p_l)
    fechas_final_final = fecha(fecha_final_p_l)
    
    fecha_titulo = fechas_inicial[0]+" " +  fechas_inicial[1] + " AL " + fechas_final[0] + " "+fechas_final[1] + " "+fechas_final[2]
    fecha_titulo = fecha_titulo.upper()

    fecha_titulo_dos = fechas_inicial_final[0]+" " +  fechas_inicial_final[1] + " AL " + fechas_final_final[0] + " "+fechas_final_final[1] + " "+fechas_final_final[2]
    fecha_titulo_dos = fecha_titulo_dos.upper()

    fecha_inicial_dia_inicial = fechas_inicial[0]
    fecha_inicial_mes_inicial = fechas_inicial[1]
    fecha_inicial_anio_inicial = fechas_inicial[2]

    fecha_inicial_dia = fechas_final[0]
    fecha_inicial_mes = fechas_final[1]
    fecha_inicial_anio = fechas_final[2]


    if fecha_inicial_anio_inicial == fecha_inicial_anio:
        fechas= str("DEL ")+str(fecha_inicial_dia_inicial) + " " +str(fecha_inicial_mes_inicial) + " hasta " + str(fecha_inicial_dia) + " de "+ str(fecha_inicial_mes) + " del " + str(fecha_inicial_anio)
    else:
        fechas= str("DEL ")+str(fecha_inicial_dia_inicial) + " " +str(fecha_inicial_mes_inicial) + " " +str(fecha_inicial_anio_inicial) + " hasta " + str(fecha_inicial_dia) + " de "+ str(fecha_inicial_mes) + " del " + str(fecha_inicial_anio)
    fechas = fechas.upper()

    esultados_spoa =  Calculo_Spoa(fecha_inicial_u_l, fecha_final_u_l, fecha_inicial_p_l, fecha_final_p_l, filtro, fechas_final[2], fechas_final_final[2], fecha_titulo, fecha_titulo_dos)

    asesinados = esultados_spoa.comparativo()


    prs = Presentation('C:/Users/nicolas.valbuena/Documents/programacion 2023/server/tipo_docker/b_f_resultados_lineas_estrategica_4_power/models/documentos/template/templete.pptx')
    slide = prs.slides.add_slide(prs.slide_layouts[2])

    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))
  
    categories = [ fechas_final_final[2], fechas_final[2]]
    series_ad =  (asesinados[2][0], asesinados[2][1])
    series_nombre = 'serie'
    charts(slide, 15, categories, series_ad, series_nombre, "ASE.")

    categories = [fechas_final_final[2], fechas_final[2]]
    series_ad =  (asesinados[3][0], asesinados[3][1])
    series_nombre = 'serie'

    charts(slide, 14, categories, series_ad, series_nombre, "HER")

    placeholder = slide.placeholders[18]
    picture = placeholder.insert_picture('C:/Users/nicolas.valbuena/Documents/programacion 2023/server/static/img/resultados/asesinados.png')

    placeholder = slide.placeholders[13]
    picture = placeholder.insert_picture('C:/Users/nicolas.valbuena/Documents/programacion 2023/server/static/img/resultados/heridos_copy.jpeg')

    tabla(slide, asesinados[0], fechas_final_final[2], fechas_final[2], 16, "ASESINADOS")
    tabla(slide, asesinados[1], fechas_final_final[2], fechas_final[2], 17, "HERIDOS")


    titulo_unidad  =  titulos_name_lineas_estrategicas_obj4(filtro, fechas_final[2])
    if len(titulo_unidad[1])>0:
        municipios=tuple(titulo_unidad[1])
    else: 
        municipios =""
    titulo= titulo_unidad[0]


    placeholder = slide.placeholders[0]
    placeholder = slide.shapes.title
    placeholder.text = titulo

    sub_titulo = "Proteger y fortalecer la Fuerza y sus capacidades estratégicas"

    if asesinados[2][0] > asesinados[2][1]:
        nombre = "tri_verde_new"
    else:
        nombre = "tri_rojo_new"
    
    direcion = 'C:/Users/nicolas.valbuena/Documents/programacion 2023/server/static/img/resultados/{}.png'.format(nombre)

    placeholder = slide.placeholders[19]
    picture = placeholder.insert_picture(direcion)
    picture.width = Inches(0.4)
    picture.height = Inches(0.4)
    picture.left = Inches(4.1)
    picture.top = Inches(6.5)
    
    cuadros_divisio(slide, (asesinados[2][2]+" %"), 4, 6.25)

    if asesinados[3][0] > asesinados[3][1]:
        nombre = "tri_verde_new"
    else:
        nombre = "tri_rojo_new"

    direcion = 'C:/Users/nicolas.valbuena/Documents/programacion 2023/server/static/img/resultados/{}.png'.format(nombre)
    placeholder = slide.placeholders[20]
    picture = placeholder.insert_picture(direcion)
    picture.width = Inches(0.4)
    picture.height = Inches(0.4)
    picture.left = Inches(8.9)
    picture.top = Inches(6.5)

    cuadros_divisio(slide, (asesinados[3][2]+" %"), 8.8, 6.25)

    placeholder = slide.placeholders[21]
    text_frame = placeholder.text_frame
    text_frame.text = sub_titulo

    placeholder = slide.placeholders[22]
    text_frame = placeholder.text_frame
    text_frame.text = fechas

    
    direcion = dirercion_archvios+str("test")+'.pptx'
    prs.save(direcion)

    direcion = dirercion_archvios+str("test")+'.pptx'

    DIRECION_3 = os.getenv('DIRECION_3')
    direcion= DIRECION_3+nombre_carpeta+str("test")+'.pptx'
    
    return [direcion, "nombre_doc"]

